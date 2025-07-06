#!/usr/bin/env python
# -*- coding: utf-8 -*-


import os
import tempfile
import argparse
import time
import multiprocessing
import subprocess
from subprocess import call, Popen, PIPE
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed
from functools import partial

from pdf2image import convert_from_path
from pptx import Presentation
from gtts import gTTS


__author__ = ['chaonan99']


## Sometimes ffmpeg is avconv
FFMPEG_NAME = 'ffmpeg'
# FFMPEG_NAME = 'avconv'


def ppt_presenter(pptx_path, pdf_path, output_path):
    with tempfile.TemporaryDirectory() as temp_path:
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        print(f"Number of PDF pages: {len(images_from_path)}")
        print(f"Number of PowerPoint slides: {len(prs.slides)}")
        assert len(images_from_path) == len(prs.slides)

        processed_slides = 0
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            print(f"Processing slide {i+1}")

            # Get notes text or use default
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                notes = slide.notes_slide.notes_text_frame.text.strip()
                print(f"Slide {i+1} has notes: {notes[:50]}...")
            else:
                notes = f"Slide {i+1}"  # Default text for slides without notes
                print(f"Slide {i+1} has no notes - using default text: '{notes}'")

            # Create TTS audio
            tts = gTTS(text=notes, lang='en')
            image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
            audio_path = os.path.join(temp_path, 'frame_{}.mp3'.format(i))

            image.save(image_path)
            tts.save(audio_path)

            ffmpeg_call(image_path, audio_path, temp_path, i)
            processed_slides += 1

        print(f"Processed {processed_slides} slides out of {len(prs.slides)} total slides")

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i)) \
                      for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)


def ffmpeg_call(image_path, audio_path, temp_path, i):
    """Original sequential FFmpeg call - kept for backward compatibility"""
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))

    # First FFmpeg call to create MP4
    cmd1 = [FFMPEG_NAME, '-loop', '1', '-y', '-i', image_path, '-i', audio_path,
            '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'aac',
            '-b:a', '192k', '-pix_fmt', 'yuv420p', '-shortest', out_path_mp4]
    print(f"Running: {' '.join(cmd1)}")
    ret1 = call(cmd1)
    print(f"Return code: {ret1}")

    # Second FFmpeg call to convert to TS
    cmd2 = [FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
            '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts]
    print(f"Running: {' '.join(cmd2)}")
    ret2 = call(cmd2)
    print(f"Return code: {ret2}")

    # Check if files were created
    print(f"MP4 exists: {os.path.exists(out_path_mp4)}")
    print(f"TS exists: {os.path.exists(out_path_ts)}")


def ffmpeg_call_parallel(image_path, audio_path, temp_path, i):
    """
    Optimized parallel FFmpeg call - uses the proven two-step approach but with better parallelization
    """
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))

    try:
        # Step 1: Create MP4 from image and audio
        cmd1 = [
            FFMPEG_NAME,
            '-loop', '1',
            '-y',
            '-i', image_path,
            '-i', audio_path,
            '-c:v', 'libx264',
            '-tune', 'stillimage',
            '-c:a', 'aac',
            '-b:a', '192k',
            '-pix_fmt', 'yuv420p',
            '-shortest',
            out_path_mp4
        ]

        result1 = subprocess.run(
            cmd1,
            capture_output=True,
            text=True,
            timeout=120  # Increased timeout to 2 minutes
        )

        if result1.returncode != 0:
            print(f"⚠️  FFmpeg MP4 creation error for slide {i+1}: {result1.stderr[:200]}")
            return False

        # Step 2: Convert MP4 to TS
        cmd2 = [
            FFMPEG_NAME,
            '-y',
            '-i', out_path_mp4,
            '-c', 'copy',
            '-bsf:v', 'h264_mp4toannexb',
            '-f', 'mpegts',
            out_path_ts
        ]

        result2 = subprocess.run(
            cmd2,
            capture_output=True,
            text=True,
            timeout=60  # TS conversion is faster
        )

        if result2.returncode != 0:
            print(f"⚠️  FFmpeg TS conversion error for slide {i+1}: {result2.stderr[:200]}")
            return False

        # Verify output file was created
        if not os.path.exists(out_path_ts):
            print(f"❌ Output TS file not created for slide {i+1}")
            return False

        # Clean up intermediate MP4 file to save space
        try:
            os.remove(out_path_mp4)
        except:
            pass  # Don't fail if cleanup fails

        return True

    except subprocess.TimeoutExpired:
        print(f"⏰ FFmpeg timeout for slide {i+1} (this may indicate a complex slide)")
        return False
    except Exception as e:
        print(f"❌ FFmpeg exception for slide {i+1}: {e}")
        return False


def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])


def process_slide_tts(slide_data):
    """Process TTS for a single slide - designed for parallel execution"""
    i, slide, temp_path = slide_data

    # Get notes text or use default
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
        notes = slide.notes_slide.notes_text_frame.text.strip()
        print(f"Slide {i+1} has notes: {notes[:50]}...")
    else:
        notes = f"Slide {i+1}"  # Default text for slides without notes
        print(f"Slide {i+1} has no notes - using default text: '{notes}'")

    # Create TTS audio
    tts = gTTS(text=notes, lang='en')
    audio_path = os.path.join(temp_path, 'frame_{}.mp3'.format(i))
    tts.save(audio_path)

    return i, audio_path, notes


def process_slide_image(image_data):
    """Process image for a single slide - designed for parallel execution"""
    i, image, temp_path = image_data

    image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
    image.save(image_path)

    return i, image_path


def process_slide_video(video_data):
    """Process video for a single slide - designed for parallel execution"""
    i, image_path, audio_path, temp_path = video_data

    print(f"🎬 Processing video for slide {i+1}")

    # Verify input files exist
    if not os.path.exists(image_path):
        print(f"❌ Image file missing for slide {i+1}: {image_path}")
        return i, False
    if not os.path.exists(audio_path):
        print(f"❌ Audio file missing for slide {i+1}: {audio_path}")
        return i, False

    success = ffmpeg_call_parallel(image_path, audio_path, temp_path, i)

    if success:
        print(f"✅ Video completed for slide {i+1}")
    else:
        print(f"❌ Video failed for slide {i+1}")

    return i, success


def process_slide_video_batch(video_batch_data):
    """Process multiple slides in a single worker - for even better efficiency"""
    batch_id, video_data_list = video_batch_data
    results = []

    print(f"🎬 Processing video batch {batch_id} with {len(video_data_list)} slides")

    for video_data in video_data_list:
        i, image_path, audio_path, temp_path = video_data
        success = ffmpeg_call_parallel(image_path, audio_path, temp_path, i)
        results.append((i, success))

        if success:
            print(f"✅ Slide {i+1} completed in batch {batch_id}")
        else:
            print(f"❌ Slide {i+1} failed in batch {batch_id}")

    return batch_id, results


def ppt_presenter_parallel(pptx_path, pdf_path, output_path, max_workers=None, use_batch_processing=False):
    """
    Parallel version of ppt_presenter with significant performance improvements.

    Args:
        pptx_path: Path to PowerPoint file
        pdf_path: Path to PDF file
        output_path: Output video path
        max_workers: Maximum number of parallel workers (None = auto-detect)
        use_batch_processing: Use batch processing for even better performance (experimental)
    """
    start_time = time.time()

    if max_workers is None:
        max_workers = min(multiprocessing.cpu_count(), 8)  # Cap at 8 to avoid overwhelming system

    print(f"Starting parallel processing with {max_workers} workers...")

    with tempfile.TemporaryDirectory() as temp_path:
        # Load images and slides
        print("Loading PDF pages and PowerPoint slides...")
        images_from_path = convert_from_path(pdf_path)
        prs = Presentation(pptx_path)
        print(f"Number of PDF pages: {len(images_from_path)}")
        print(f"Number of PowerPoint slides: {len(prs.slides)}")
        assert len(images_from_path) == len(prs.slides)

        # Phase 1: Parallel TTS generation (I/O bound - use ThreadPoolExecutor)
        print("\nPhase 1: Generating TTS audio files in parallel...")
        tts_start = time.time()

        slide_tts_data = [(i, slide, temp_path) for i, slide in enumerate(prs.slides)]
        audio_paths = {}

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            tts_futures = {executor.submit(process_slide_tts, data): data[0] for data in slide_tts_data}

            for future in as_completed(tts_futures):
                i, audio_path, notes = future.result()
                audio_paths[i] = audio_path
                print(f"✓ TTS completed for slide {i+1}")

        tts_time = time.time() - tts_start
        print(f"TTS generation completed in {tts_time:.2f} seconds")

        # Phase 2: Parallel image processing (I/O bound - use ThreadPoolExecutor)
        print("\nPhase 2: Processing images in parallel...")
        image_start = time.time()

        image_data = [(i, image, temp_path) for i, image in enumerate(images_from_path)]
        image_paths = {}

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            image_futures = {executor.submit(process_slide_image, data): data[0] for data in image_data}

            for future in as_completed(image_futures):
                i, image_path = future.result()
                image_paths[i] = image_path
                print(f"✓ Image processed for slide {i+1}")

        image_time = time.time() - image_start
        print(f"Image processing completed in {image_time:.2f} seconds")

        # Phase 3: Parallel video generation (CPU bound - use ProcessPoolExecutor)
        print("\nPhase 3: Generating videos in parallel...")
        video_start = time.time()

        video_data = [(i, image_paths[i], audio_paths[i], temp_path)
                     for i in range(len(prs.slides))]

        # Optimize worker count for FFmpeg operations
        # Use conservative worker count to avoid system overload
        video_workers = min(max_workers // 2, multiprocessing.cpu_count() // 2, 4)
        video_workers = max(video_workers, 1)  # Ensure at least 1 worker
        print(f"Using {video_workers} workers for video processing...")

        # Option 1: Individual slide processing (good for mixed slide complexity)
        failed_slides = []
        with ProcessPoolExecutor(max_workers=video_workers) as executor:
            video_futures = {executor.submit(process_slide_video, data): data[0] for data in video_data}

            completed_videos = 0
            for future in as_completed(video_futures):
                i, success = future.result()
                completed_videos += 1
                if not success:
                    failed_slides.append(i)
                print(f"✓ Video processed for slide {i+1} ({completed_videos}/{len(prs.slides)})")

        # Report any failures
        if failed_slides:
            print(f"⚠️  {len(failed_slides)} slides failed: {[i+1 for i in failed_slides]}")
        else:
            print(f"🎉 All {len(prs.slides)} videos generated successfully!")

        video_time = time.time() - video_start
        print(f"Video generation completed in {video_time:.2f} seconds")

        # Phase 4: Concatenate all videos
        print("\nPhase 4: Concatenating final video...")
        concat_start = time.time()

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i))
                     for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)

        concat_time = time.time() - concat_start
        total_time = time.time() - start_time

        print(f"Final concatenation completed in {concat_time:.2f} seconds")
        print(f"\n🎉 Total processing time: {total_time:.2f} seconds")
        print(f"   - TTS generation: {tts_time:.2f}s")
        print(f"   - Image processing: {image_time:.2f}s")
        print(f"   - Video generation: {video_time:.2f}s")
        print(f"   - Final concatenation: {concat_time:.2f}s")
        print(f"✅ Parallel video generation complete: {output_path}")


def compare_performance(pptx_path, pdf_path, output_path, workers=None):
    """Compare performance between sequential and parallel processing"""
    print("🔬 Performance Comparison Mode")
    print("=" * 50)

    # Test sequential processing
    print("\n1️⃣  Testing Sequential Processing...")
    sequential_output = output_path.replace('.mp4', '_sequential.mp4')
    start_time = time.time()
    ppt_presenter(pptx_path, pdf_path, sequential_output)
    sequential_time = time.time() - start_time

    # Test parallel processing
    print("\n2️⃣  Testing Parallel Processing...")
    parallel_output = output_path.replace('.mp4', '_parallel.mp4')
    start_time = time.time()
    ppt_presenter_parallel(pptx_path, pdf_path, parallel_output, workers)
    parallel_time = time.time() - start_time

    # Results
    speedup = sequential_time / parallel_time if parallel_time > 0 else 0
    print("\n📊 Performance Results:")
    print("=" * 50)
    print(f"Sequential time: {sequential_time:.2f} seconds")
    print(f"Parallel time:   {parallel_time:.2f} seconds")
    print(f"Speedup:         {speedup:.2f}x faster")
    print(f"Time saved:      {sequential_time - parallel_time:.2f} seconds")
    print(f"\n✅ Both videos generated:")
    print(f"   Sequential: {sequential_output}")
    print(f"   Parallel:   {parallel_output}")


def main():
    parser = argparse.ArgumentParser(
        description='PPT Presenter - Convert PowerPoint to video with narration',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Basic usage (parallel mode by default)
  python ppt_presenter.py --pptx slides.pptx --pdf slides.pdf -o output.mp4

  # Use sequential mode
  python ppt_presenter.py --pptx slides.pptx --pdf slides.pdf -o output.mp4 --sequential

  # Turbo mode for maximum speed (uses more CPU/memory)
  python ppt_presenter.py --pptx slides.pptx --pdf slides.pdf -o output.mp4 --turbo

  # Specify number of workers
  python ppt_presenter.py --pptx slides.pptx --pdf slides.pdf -o output.mp4 --workers 4

  # Compare performance
  python ppt_presenter.py --pptx slides.pptx --pdf slides.pdf -o output.mp4 --compare
        """)

    parser.add_argument('--pptx', required=True, help='input PowerPoint (.pptx) file path')
    parser.add_argument('--pdf', required=True, help='input PDF file path')
    parser.add_argument('-o', '--output', required=True, help='output video file path')
    parser.add_argument('--parallel', action='store_true', default=True,
                       help='use parallel processing (default: True)')
    parser.add_argument('--sequential', action='store_true',
                       help='use sequential processing (overrides --parallel)')
    parser.add_argument('--workers', type=int, default=None,
                       help='number of parallel workers (default: auto-detect based on CPU cores)')
    parser.add_argument('--turbo', action='store_true',
                       help='enable turbo mode with maximum parallelization (uses more CPU/memory)')
    parser.add_argument('--compare', action='store_true',
                       help='compare performance between sequential and parallel modes')

    args = parser.parse_args()

    # Performance comparison mode
    if args.compare:
        compare_performance(args.pptx, args.pdf, args.output, args.workers)
        return

    # Determine processing mode
    use_parallel = args.parallel and not args.sequential

    if use_parallel:
        if args.turbo:
            print("🚀💨 Using TURBO parallel processing mode (maximum speed)")
            # In turbo mode, use more aggressive worker settings
            turbo_workers = args.workers if args.workers else min(multiprocessing.cpu_count() * 2, 16)
            ppt_presenter_parallel(args.pptx, args.pdf, args.output, turbo_workers)
        else:
            print("🚀 Using parallel processing mode")
            ppt_presenter_parallel(args.pptx, args.pdf, args.output, args.workers)
    else:
        print("🐌 Using sequential processing mode")
        ppt_presenter(args.pptx, args.pdf, args.output)


if __name__ == '__main__':
    main()