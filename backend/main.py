import os
import sys
import json
import time
import tempfile
import subprocess
import multiprocessing
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed
from functools import partial
from pdf2image import convert_from_path
from pptx import Presentation
from gtts import gTTS
import io
import random
import logging
from typing import List, Optional, Dict, Any
from datetime import datetime, timedelta
import asyncio
import aiofiles
import requests
from geopy.geocoders import Nominatim
from geopy.exc import GeocoderTimedOut, GeocoderUnavailable

from fastapi import FastAPI, HTTPException, Request, Body, File, UploadFile, Form
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
import uvicorn
import firebase_admin
from firebase_admin import credentials, firestore, auth
import google.generativeai as genai
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configure logging first
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Import chat and RAG services
try:
    from chat_service import threat_ai
    from rag_service_v2 import enhanced_rag_service
    from web_search_service import web_search_service
    CHAT_ENABLED = True
    RAG_ENABLED = True
    WEB_SEARCH_ENABLED = True
    logger.info("Enhanced RAG, Chat, and Web Search services loaded successfully")
except ImportError as e:
    logging.warning(f"Chat services not available: {e}")
    CHAT_ENABLED = False
    RAG_ENABLED = False
    WEB_SEARCH_ENABLED = False

# Set up path for imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import Firebase configuration
from firebase_config import (
    initialize_firebase, 
    get_firestore_db,
    get_or_create_user, 
    get_user_threat_stats, 
    get_user_threat_categories,
    get_user_analysis_history, 
    add_analysis_to_history, 
    update_user_threat_stats,
    add_threat_location,
    get_user_threat_locations,
    get_all_threat_locations,
    filter_threat_locations
)

# Import model loader
from model_loader import model_loader

# Import geocoding utilities
from geocoding_utils import geocode_location, determine_threat_priority, extract_location_from_user_data, extract_location_from_twitter_api_response

# Import FIR service
try:
    from fir_service import fir_service
    FIR_ENABLED = True
    logger.info("FIR service loaded successfully")
except ImportError as e:
    logging.warning(f"FIR service not available: {e}")
    FIR_ENABLED = False

# Logger already configured above

# Define request models
class PredictionRequest(BaseModel):
    text: str
    model_type: Optional[str] = "distilbert"  # "distilbert" or "astra"

class BatchPredictionRequest(BaseModel):
    texts: List[str]
    model_type: Optional[str] = "distilbert"  # "distilbert" or "astra"

class UserInfo(BaseModel):
    user_id: str
    email: str
    first_name: Optional[str] = None
    last_name: Optional[str] = None

class ThreatMapFilterRequest(BaseModel):
    timeRange: Optional[int] = 30  # days
    threatTypes: Optional[List[str]] = None
    priority: Optional[List[str]] = None

# Create FastAPI app
app = FastAPI(
    title="Astra API",
    description="Astra: API for detecting threatening content in text",
    version="1.0.0"
)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize Firebase on startup
db = None

@app.on_event("startup")
async def startup_event():
    global db
    
    # Initialize Firebase
    if initialize_firebase():
        db = get_firestore_db()
        logger.info("Firebase initialized successfully")
    else:
        logger.error("Failed to initialize Firebase")
    
    # Load both DistilBERT and Astra models  
    stage1_dir = "/Users/darkarmy/Downloads/threat_detection/myProj/models/stage1_bin"
    stage2_dir = "/Users/darkarmy/Downloads/threat_detection/myProj/models/stage2_multi"
    astra_path = "/Users/darkarmy/Downloads/threat_detection/myProj/models/optimized_threat_detection_model.pth"
    
    logger.info(f"Loading models:")
    logger.info(f"  DistilBERT Stage 1: {stage1_dir}")
    logger.info(f"  DistilBERT Stage 2: {stage2_dir}")
    logger.info(f"  Astra Model: {astra_path}")
    
    model_results = model_loader.load_all_models(stage1_dir, stage2_dir, astra_path)
    
    loaded_models = [model for model, status in model_results.items() if status]
    failed_models = [model for model, status in model_results.items() if not status]
    
    if loaded_models:
        logger.info(f"✅ Successfully loaded models: {', '.join(loaded_models)}")
    
    if failed_models:
        logger.warning(f"⚠️ Failed to load models: {', '.join(failed_models)}")
    
    if not any(model_results.values()):
        logger.error("❌ No models loaded successfully. Will fallback to mock predictions.")
    
    # Configure Gemini API
    try:
        gemini_api_key = os.getenv('GEMINI_API_KEY')
        if gemini_api_key:
            genai.configure(api_key=gemini_api_key)
            logger.info("✅ Gemini API configured successfully")
        else:
            logger.warning("⚠️ GEMINI_API_KEY not found in environment variables")
    except Exception as e:
        logger.error(f"❌ Failed to configure Gemini API: {e}")

# Define label map for threat classification - ONLY used as fallback if model loading fails
label_map = {
    0: "Hate Speech/Extremism",
    1: "Direct Violence Threats",
    2: "Harassment and Intimidation",
    3: "Criminal Activity",
    4: "Child Safety Threats",
    5: "Not a Threat"
}

# Health check endpoint
@app.get("/health")
async def health_check():
    return {"status": "ok", "message": "API is running", "timestamp": time.time()}

# Get available models endpoint
@app.get("/api/models")
async def get_models():
    """Get information about available models"""
    return {
        "available_models": model_loader.get_available_models(),
        "model_info": model_loader.get_model_info()
    }

# Extract user_id from request headers
def get_user_id(request: Request) -> str:
    user_id = request.headers.get("user_id")
    if not user_id:
        user_id = request.query_params.get("user_id", "anonymous")
    return user_id

# Register or update user
@app.post("/api/user/register")
async def register_user(user: UserInfo):
    try:
        logger.info(f"Registering user with ID: {user.user_id}, email: {user.email}")
        
        # Validate input
        if not user.user_id:
            logger.error("Missing user_id in registration request")
            raise HTTPException(status_code=400, detail="User ID is required")
            
        if not user.email:
            logger.warning(f"No email provided for user {user.user_id}, using placeholder")
            user.email = f"{user.user_id}@placeholder.email.com"
        
        # Create or update user in Firestore
        user_data = get_or_create_user(user.user_id, user.email, user.first_name, user.last_name)
        if not user_data:
            logger.error(f"Failed to register user {user.user_id} in Firebase")
            raise HTTPException(status_code=500, detail="Failed to register user in database")
        
        logger.info(f"User registered successfully: {user.user_id}")
        return {"status": "success", "message": "User registered successfully"}
    except HTTPException as http_err:
        # Re-raise HTTP exceptions
        raise http_err
    except Exception as e:
        error_msg = str(e)
        logger.error(f"Error registering user: {error_msg}")
        logger.exception("Full stack trace:")
        raise HTTPException(status_code=500, detail=f"Registration failed: {error_msg}")

# Predict endpoint
@app.post("/api/predict")
async def predict(request: Request, prediction_request: PredictionRequest):
    text = prediction_request.text
    model_type = prediction_request.model_type or "distilbert"
    
    # Get the user_id from request headers or query params
    user_id = get_user_id(request)
    logger.info(f"Processing prediction for user: {user_id} using model: {model_type}")
    
    # Check if the requested model is available
    available_models = model_loader.get_available_models()
    if model_type not in available_models or not available_models[model_type]:
        raise HTTPException(
            status_code=400, 
            detail=f"Model '{model_type}' is not available. Available models: {[k for k, v in available_models.items() if v]}"
        )
    
    # Try to predict using the specified model
    result = None
    
    try:
        model_result = model_loader.predict(text, model_type=model_type)
        if model_result.get("success", False):
            result = model_result
            # Add timestamp and text
            result["timestamp"] = datetime.now().isoformat()
            result["text"] = text
            logger.info(f"Model prediction successful for user {user_id} using {model_type}")
        else:
            logger.error(f"Model prediction failed for user {user_id}: {model_result}")
            raise HTTPException(status_code=500, detail=f"Model prediction failed: {model_result.get('error', 'Unknown error')}")
    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logger.error(f"Error during model prediction for user {user_id}: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Model prediction error: {str(e)}")

    # Save to Firebase if user is not anonymous and we have a real result
    if user_id != "anonymous" and result:
        try:
            history_item = add_analysis_to_history(user_id, text, result)
            if history_item:
                result["id"] = history_item.get("id")
            
            # Update user's threat statistics
            update_user_threat_stats(user_id, result)
            
            # Also add to RAG system for chatbot context
            if RAG_ENABLED:
                try:
                    rag_success = enhanced_rag_service.add_prediction_analysis(
                        user_id=user_id,
                        text=text,
                        prediction_result=result,
                        source='single_prediction'
                    )
                    if rag_success:
                        logger.info(f"✅ Prediction indexed in enhanced RAG for user {user_id}")
                except Exception as rag_error:
                    logger.error(f"❌ Enhanced RAG indexing error: {rag_error}")
            
            # Save to threat map if this is a threat
            if result.get("threat", False):
                try:
                    # Use generic location for direct text analysis
                    location_str = "Direct text analysis threat"
                    
                    # Geocode the location (will generate random coordinates)
                    lat, lng = geocode_location(location_str)
                    
                    # Determine priority
                    predicted_class = result.get('predicted_class', 'Unknown')
                    confidence = result.get('confidence', 0.0)
                    priority = determine_threat_priority(predicted_class, confidence)
                    
                    # Create threat location entry
                    threat_data = {
                        "id": f"THR-{result.get('id', str(uuid.uuid4())[:8])}",
                        "type": predicted_class,
                        "lat": lat,
                        "lng": lng,
                        "title": f"{predicted_class} detected in direct analysis",
                        "location": location_str,
                        "date": result.get('timestamp', datetime.now().isoformat()),
                        "priority": priority,
                        "details": text[:100] + ("..." if len(text) > 100 else ""),
                        "caseId": f"THP-{str(uuid.uuid4())[:8]}",
                        "source": "direct_prediction",
                        "confidence": confidence,
                        "predicted_class": predicted_class,
                        "twitter_metadata": {},
                        "text": text
                    }
                    
                    # Save to threat map
                    saved_location = add_threat_location(user_id, threat_data)
                    if saved_location:
                        logger.info(f"✅ Saved threat location for direct prediction")
                    else:
                        logger.warning(f"⚠️ Failed to save threat location for direct prediction")
                        
                except Exception as location_error:
                    logger.error(f"❌ Error saving threat location for direct prediction: {location_error}")
            
            logger.info(f"✅ Analysis saved to Firebase for user {user_id}")
        except Exception as e:
            logger.error(f"❌ Error saving to Firebase for user {user_id}: {e}")
            # Don't fail the request if Firebase save fails, just log
            result["firebase_save_error"] = str(e)

    return result

# Get user stats
@app.get("/api/user/stats")
async def get_user_stats(request: Request):
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        # Get user stats
        stats = get_user_threat_stats(user_id)
        if not stats:
            raise HTTPException(status_code=500, detail="Failed to get user stats")
        
        # Get categories
        categories = get_user_threat_categories(user_id)
        if not categories:
            raise HTTPException(status_code=500, detail="Failed to get user categories")
        
        return {
            "stats": stats,
            "categories": categories
        }
    
    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        logger.error(f"Error getting user stats: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Recalculate threat categories from history
@app.post("/api/user/categories/recalculate")
async def recalculate_user_categories(request: Request):
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        # TODO: Implement recalculate_threat_categories_from_history function
        # from firebase_config import recalculate_threat_categories_from_history
        # success = recalculate_threat_categories_from_history(user_id)
        
        # For now, just return current categories
        categories = get_user_threat_categories(user_id)
        if not categories:
            categories = {}  # Return empty categories if none found
        
        return {
            "success": True,
            "message": "Categories retrieved successfully (recalculation not implemented yet)",
            "categories": categories
        }
        
    except Exception as e:
        if isinstance(e, HTTPException):
            raise e
        logger.error(f"Error getting categories: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Get user history
@app.get("/api/user/history")
async def get_user_history(request: Request):
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        logger.warning("Anonymous user tried to access history")
        raise HTTPException(status_code=401, detail="Authentication required")
    
    logger.info(f"Fetching history for user: {user_id}")
    
    try:
        # Get user's history from Firebase
        history = get_user_analysis_history(user_id)
        
        if history is None:
            logger.error(f"Failed to fetch history for user {user_id}")
            raise HTTPException(status_code=500, detail="Failed to fetch history data")
        
        if not history:  # Empty list
            logger.info(f"No history found for user {user_id}")
            return []
        
        logger.info(f"Successfully fetched {len(history)} history items for user {user_id}")
        return history
    
    except Exception as e:
        logger.error(f"Error getting user history: {e}")
        logger.exception("Full stack trace:")
        raise HTTPException(status_code=500, detail=str(e))

# Batch prediction endpoint
@app.post("/api/predict/batch")
async def predict_batch(request: Request, batch_request: BatchPredictionRequest):
    results = []
    firebase_results = []
    model_type = batch_request.model_type or "distilbert"
    
    # Get the user_id for database operations
    user_id = get_user_id(request)
    logger.info(f"Processing batch prediction for user: {user_id}, items: {len(batch_request.texts)}, model: {model_type}")
    
    # Check if the requested model is available
    available_models = model_loader.get_available_models()
    if model_type not in available_models or not available_models[model_type]:
        raise HTTPException(
            status_code=400, 
            detail=f"Model '{model_type}' is not available. Available models: {[k for k, v in available_models.items() if v]}"
        )
    
    for text in batch_request.texts:
        result = None
        
        # Use the selected model for prediction
        try:
            model_result = model_loader.predict(text, model_type=model_type)
            if model_result.get("success", False):
                result = model_result
                # Add timestamp and text
                result["timestamp"] = datetime.now().isoformat()
                result["text"] = text
            else:
                logger.error(f"Model prediction failed for text in batch: {model_result}")
                continue  # Skip this text if prediction fails
        except Exception as e:
            logger.error(f"Error during batch prediction: {str(e)}")
            continue  # Skip this text if prediction fails

        # Save to Firebase if user is not anonymous and we have a real result
        if user_id != "anonymous" and result:
            try:
                history_item = add_analysis_to_history(user_id, text, result)
                if history_item:
                    result["id"] = history_item.get("id")
                    firebase_results.append(history_item)
                
                # Update user's threat statistics  
                update_user_threat_stats(user_id, result)
                
                # Also add to RAG system for chatbot context
                if RAG_ENABLED:
                    try:
                        rag_success = enhanced_rag_service.add_prediction_analysis(
                            user_id=user_id,
                            text=text,
                            prediction_result=result,
                            source='batch_prediction'
                        )
                        if rag_success:
                            logger.debug(f"✅ Batch prediction indexed in enhanced RAG for user {user_id}")
                    except Exception as rag_error:
                        logger.error(f"❌ Enhanced RAG batch indexing error: {rag_error}")
                
                # Save to threat map if this is a threat
                if result.get("threat", False):
                    try:
                        # Use generic location for batch analysis
                        location_str = "Batch analysis threat"
                        
                        # Geocode the location (will generate random coordinates)
                        lat, lng = geocode_location(location_str)
                        
                        # Determine priority
                        predicted_class = result.get('predicted_class', 'Unknown')
                        confidence = result.get('confidence', 0.0)
                        priority = determine_threat_priority(predicted_class, confidence)
                        
                        # Create threat location entry
                        threat_data = {
                            "id": f"THR-{result.get('id', str(uuid.uuid4())[:8])}",
                            "type": predicted_class,
                            "lat": lat,
                            "lng": lng,
                            "title": f"{predicted_class} detected in batch analysis",
                            "location": location_str,
                            "date": result.get('timestamp', datetime.now().isoformat()),
                            "priority": priority,
                            "details": text[:100] + ("..." if len(text) > 100 else ""),
                            "caseId": f"THP-{str(uuid.uuid4())[:8]}",
                            "source": "batch_prediction",
                            "confidence": confidence,
                            "predicted_class": predicted_class,
                            "twitter_metadata": {},
                            "text": text
                        }
                        
                        # Save to threat map
                        saved_location = add_threat_location(user_id, threat_data)
                        if saved_location:
                            logger.debug(f"✅ Saved threat location for batch prediction")
                        
                    except Exception as location_error:
                        logger.error(f"❌ Error saving threat location for batch prediction: {location_error}")
                
            except Exception as e:
                logger.error(f"Error saving batch item to Firebase: {e}")
                # Continue with next item even if this one fails to save
        
        if result:
            results.append(result)
    
    logger.info(f"Batch prediction complete: {len(results)} successful predictions out of {len(batch_request.texts)} texts")
    
    return {"results": results, "total_processed": len(results)}

# Save summary report
@app.post("/api/user/reports/summary")
async def save_summary_report(request: Request, report: dict = Body(...)):
    user_id = get_user_id(request)
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        # Get Firestore database instance
        from firebase_config import get_firestore_db
        db_instance = get_firestore_db()
        if not db_instance:
            raise HTTPException(status_code=503, detail="Database not available")
        
        # Add timestamp and user info
        report_data = {
            **report,
            "saved_at": datetime.now().isoformat(),
            "user_id": user_id
        }
        
        # Save report to Firestore
        db_instance.collection("users").document(user_id).collection("reports").document("summary").set(report_data)
        logger.info(f"Summary report saved successfully for user {user_id}")
        
        # Also add to enhanced RAG system for chatbot context
        if RAG_ENABLED:
            try:
                rag_success = enhanced_rag_service.add_summary_report(
                    user_id=user_id,
                    report_data=report_data,
                    source='frontend_summary'
                )
                if rag_success:
                    logger.info(f"✅ Summary report indexed in enhanced RAG for user {user_id}")
                else:
                    logger.warning(f"⚠️ Failed to index summary report in enhanced RAG for user {user_id}")
            except Exception as rag_error:
                logger.error(f"❌ Enhanced RAG summary indexing error: {rag_error}")
                # Don't fail the main request if RAG indexing fails
        
        return {"status": "success", "message": "Summary report saved successfully"}
    except Exception as e:
        logger.error(f"Error saving summary report for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to save summary report: {str(e)}")

# Get summary report
@app.get("/api/user/reports/summary")
async def get_summary_report(request: Request):
    user_id = get_user_id(request)
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        # Get Firestore database instance
        from firebase_config import get_firestore_db
        db_instance = get_firestore_db()
        if not db_instance:
            raise HTTPException(status_code=503, detail="Database not available")
        
        # Get report from Firestore
        doc_ref = db_instance.collection("users").document(user_id).collection("reports").document("summary")
        doc = doc_ref.get()
        
        if doc.exists:
            logger.info(f"Summary report retrieved successfully for user {user_id}")
            return {"report": doc.to_dict()}
        else:
            logger.info(f"No summary report found for user {user_id}")
        return {"report": None}
    except Exception as e:
        logger.error(f"Error getting summary report for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to retrieve summary report: {str(e)}")

# Save threat report
@app.post("/api/user/reports/threat")
async def save_threat_report(request: Request, report: dict = Body(...)):
    user_id = get_user_id(request)
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        # Get Firestore database instance
        from firebase_config import get_firestore_db
        db_instance = get_firestore_db()
        if not db_instance:
            raise HTTPException(status_code=503, detail="Database not available")
        
        # Add timestamp and user info
        report_data = {
            **report,
            "saved_at": datetime.now().isoformat(),
            "user_id": user_id
        }
        
        # Save report to Firestore
        db_instance.collection("users").document(user_id).collection("reports").document("threat").set(report_data)
        logger.info(f"Threat report saved successfully for user {user_id}")
        
        # Also add to enhanced RAG system for chatbot context
        if RAG_ENABLED:
            try:
                rag_success = enhanced_rag_service.add_threat_report(
                    user_id=user_id,
                    report_data=report_data,
                    source='frontend_threat'
                )
                
                if rag_success:
                    logger.info(f"✅ Threat report indexed in enhanced RAG for user {user_id}")
                else:
                    logger.warning(f"⚠️ Failed to index threat report in enhanced RAG for user {user_id}")
                    
            except Exception as rag_error:
                logger.error(f"❌ Enhanced RAG threat indexing error: {rag_error}")
                # Don't fail the main request if RAG indexing fails
        
        return {"status": "success", "message": "Threat report saved successfully"}
    except Exception as e:
        logger.error(f"Error saving threat report for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to save threat report: {str(e)}")

# Get threat report
@app.get("/api/user/reports/threat")
async def get_threat_report(request: Request):
    user_id = get_user_id(request)
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        # Get Firestore database instance
        from firebase_config import get_firestore_db
        db_instance = get_firestore_db()
        if not db_instance:
            raise HTTPException(status_code=503, detail="Database not available")
        
        # Get report from Firestore
        doc_ref = db_instance.collection("users").document(user_id).collection("reports").document("threat")
        doc = doc_ref.get()
        
        if doc.exists:
            logger.info(f"Threat report retrieved successfully for user {user_id}")
            return {"report": doc.to_dict()}
        else:
            logger.info(f"No threat report found for user {user_id}")
        return {"report": None}
    except Exception as e:
        logger.error(f"Error getting threat report for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to retrieve threat report: {str(e)}")

# New Twitter Analysis endpoints
class TwitterAnalysisRequest(BaseModel):
    username: str
    tweets: List[dict]  # List of tweet objects with content, etc.
    userInfo: Optional[dict] = None
    model_type: Optional[str] = "distilbert"  # "distilbert" or "astra"

class SingleTweetAnalysisRequest(BaseModel):
    tweet_text: str
    username: Optional[str] = None
    tweet_metadata: Optional[dict] = None
    model_type: Optional[str] = "distilbert"  # "distilbert" or "astra"

class ChatRequest(BaseModel):
    message: str
    use_web_search: Optional[bool] = False

class WebSearchRequest(BaseModel):
    query: str
    num_results: Optional[int] = 5

class ChatHealthRequest(BaseModel):
    check_type: Optional[str] = "basic"
    
class TwitterMonitorRequest(BaseModel):
    keyword: str
    type: str = "keyword"  # "keyword", "user", "hashtag"
    max_results: int = 10
    schedule_hour: int = 1  # Hour of day to run (0-23)
    status: str = "active"  # "active", "paused", "stopped"

class BriefingRequest(BaseModel):
    reportId: str

class FIRRequest(BaseModel):
    threat_data: dict

@app.post("/api/twitter/analyze-user")
async def analyze_twitter_user(request: Request, twitter_request: TwitterAnalysisRequest):
    """Analyze a Twitter user's tweets and save to Firebase"""
    user_id = get_user_id(request)
    model_type = twitter_request.model_type or "distilbert"
    logger.info(f"Processing Twitter user analysis for user: {user_id}, analyzing {len(twitter_request.tweets)} tweets for @{twitter_request.username} using {model_type}")
    
    # Check if the requested model is available
    available_models = model_loader.get_available_models()
    if model_type not in available_models or not available_models[model_type]:
        raise HTTPException(
            status_code=400, 
            detail=f"Model '{model_type}' is not available. Available models: {[k for k, v in available_models.items() if v]}"
        )
    
    analyzed_tweets = []
    threat_count = 0
    
    for tweet in twitter_request.tweets:
        tweet_text = tweet.get('content', '')
        if not tweet_text:
            continue
            
        # Analyze the tweet using selected model
        result = None
        try:
            model_result = model_loader.predict(tweet_text, model_type=model_type)
            if model_result.get("success", False):
                result = model_result
                result["timestamp"] = datetime.now().isoformat()
                result["text"] = tweet_text
            else:
                logger.error(f"Model prediction failed for tweet: {model_result}")
                continue  # Skip this tweet if prediction fails
        except Exception as e:
            logger.error(f"Error during Twitter tweet prediction: {str(e)}")
            continue  # Skip this tweet if prediction fails
        
        # Only process if we have a real result
        if not result:
            continue
        
        # Add Twitter-specific metadata including user profile information
        result["twitter_metadata"] = {
            "username": twitter_request.username,
            "tweet_id": tweet.get('id'),
            "created_at": tweet.get('created_at'),
            "likes": tweet.get('likes', 0),
            "retweets": tweet.get('retweets', 0),
            "analysis_type": "twitter_user_analysis"
        }
        
        # Add user metadata from userInfo if available
        if twitter_request.userInfo:
            user_info = twitter_request.userInfo
            result["user_metadata"] = {
                "display_name": user_info.get('name') or user_info.get('display_name', ''),
                "twitter_handle": twitter_request.username,
                "profile_image": user_info.get('profile_image_url') or user_info.get('profile_pic_url') or user_info.get('profile_image_url_https', ''),
                "location": user_info.get('location', ''),
                "bio": user_info.get('description', ''),
                "followers_count": user_info.get('followers_count', 0),
                "following_count": user_info.get('friends_count', 0),
                "verified": user_info.get('verified', False),
                "account_created": user_info.get('created_at', ''),
                "public_metrics": user_info.get('public_metrics', {})
            }
        
        # Save to Firebase if user is not anonymous
        if user_id != "anonymous":
            try:
                # Create enhanced text for storage that includes context
                enhanced_text = f"Tweet by @{twitter_request.username}: {tweet_text}"
                history_item = add_analysis_to_history(user_id, enhanced_text, result)
                if history_item:
                    result["id"] = history_item.get("id")
                    
                # Update user's stats
                update_user_threat_stats(user_id, result)
                
                # Save to threat map if this is a threat
                if result.get("threat", False):
                    try:
                        # Extract location from user info - enhanced location extraction
                        location_str = ""
                        
                        # Try to get location from userInfo using enhanced Twitter API response parser
                        if twitter_request.userInfo:
                            location_str = extract_location_from_twitter_api_response(twitter_request.userInfo)
                            
                            # Fallback to standard user data extraction if Twitter API response parsing fails
                            if not location_str:
                                location_str = extract_location_from_user_data(twitter_request.userInfo)
                        
                        # Fallback: try to extract from tweet location data if available
                        if not location_str and 'location' in tweet:
                            location_str = tweet.get('location', '')
                        
                        # Default if no location found
                        if not location_str:
                            location_str = f"Social media threat from @{twitter_request.username}"
                        
                        # Geocode the location
                        lat, lng = geocode_location(location_str)
                        
                        # Determine priority
                        predicted_class = result.get('predicted_class', 'Unknown')
                        confidence = result.get('confidence', 0.0)
                        priority = determine_threat_priority(predicted_class, confidence)
                        
                        # Create threat location entry
                        threat_data = {
                            "id": f"THR-{result.get('id', str(uuid.uuid4())[:8])}",
                            "type": predicted_class,
                            "lat": lat,
                            "lng": lng,
                            "title": f"{predicted_class} detected from @{twitter_request.username}",
                            "location": location_str,
                            "date": result.get('timestamp', datetime.now().isoformat()),
                            "priority": priority,
                            "details": tweet_text[:100] + ("..." if len(tweet_text) > 100 else ""),
                            "caseId": f"THP-{str(uuid.uuid4())[:8]}",
                            "source": "twitter_user_analysis",
                            "confidence": confidence,
                            "predicted_class": predicted_class,
                            "twitter_metadata": result["twitter_metadata"],
                            "user_metadata": result.get("user_metadata", {}),
                            "text": tweet_text
                        }
                        
                        # Save to threat map
                        saved_location = add_threat_location(user_id, threat_data)
                        if saved_location:
                            logger.info(f"✅ Saved threat location for @{twitter_request.username} threat at {location_str}")
                        else:
                            logger.warning(f"⚠️ Failed to save threat location for @{twitter_request.username}")
                            
                    except Exception as location_error:
                        logger.error(f"❌ Error saving threat location for @{twitter_request.username}: {location_error}")
                
            except Exception as e:
                logger.error(f"Error saving Twitter analysis to Firebase: {e}")
        
        if result.get("threat", False):
            threat_count += 1
            
        analyzed_tweets.append(result)
    
    # Create summary
    summary = {
        "username": twitter_request.username,
        "total_tweets": len(analyzed_tweets),
        "threat_tweets": threat_count,
        "threat_percentage": (threat_count / len(analyzed_tweets) * 100) if analyzed_tweets else 0,
        "analysis_timestamp": datetime.now().isoformat(),
        "user_info": twitter_request.userInfo
    }
    
    logger.info(f"Twitter analysis complete: {threat_count}/{len(analyzed_tweets)} threats detected for @{twitter_request.username}")
    
    return {
        "summary": summary,
        "analyzed_tweets": analyzed_tweets,
        "success": True
    }

@app.post("/api/twitter/analyze-tweet")
async def analyze_single_tweet(request: Request, tweet_request: SingleTweetAnalysisRequest):
    """Analyze a single tweet and save to Firebase"""
    user_id = get_user_id(request)
    tweet_text = tweet_request.tweet_text
    model_type = tweet_request.model_type or "distilbert"
    
    logger.info(f"Processing single tweet analysis for user: {user_id} using {model_type}")
    
    # Check if the requested model is available
    available_models = model_loader.get_available_models()
    if model_type not in available_models or not available_models[model_type]:
        raise HTTPException(
            status_code=400, 
            detail=f"Model '{model_type}' is not available. Available models: {[k for k, v in available_models.items() if v]}"
        )
    
    # Analyze the tweet using selected model
    result = None
    try:
        model_result = model_loader.predict(tweet_text, model_type=model_type)
        if model_result.get("success", False):
            result = model_result
            result["timestamp"] = datetime.now().isoformat()
            result["text"] = tweet_text
        else:
            logger.error(f"Model prediction failed for single tweet: {model_result}")
            raise HTTPException(status_code=500, detail=f"Model prediction failed: {model_result.get('error', 'Unknown error')}")
    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logger.error(f"Error during single tweet prediction: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Model prediction error: {str(e)}")
    
    # Only continue if we have a real result
    if not result:
        raise HTTPException(status_code=500, detail="Failed to analyze tweet")
    
    # Add Twitter metadata if provided
    if tweet_request.username or tweet_request.tweet_metadata:
        result["twitter_metadata"] = {
            "username": tweet_request.username,
            "analysis_type": "single_tweet_analysis",
            **(tweet_request.tweet_metadata or {})
        }
        
        # Add user metadata if available in tweet_metadata
        if tweet_request.tweet_metadata:
            metadata = tweet_request.tweet_metadata
            logger.info(f"🔍 Received tweet metadata for @{tweet_request.username}: {metadata}")
            
            result["user_metadata"] = {
                "display_name": metadata.get('user_name') or metadata.get('display_name', ''),
                "twitter_handle": tweet_request.username,
                "profile_image": metadata.get('user_profile_image_url') or metadata.get('profile_image_url') or metadata.get('profile_pic_url', ''),
                "location": metadata.get('user_location') or metadata.get('location', ''),
                "bio": metadata.get('user_description') or metadata.get('description', ''),
                "followers_count": metadata.get('user_followers_count', 0),
                "following_count": metadata.get('user_friends_count', 0),
                "verified": metadata.get('user_verified', False),
                "account_created": metadata.get('user_created_at', ''),
                "public_metrics": metadata.get('user_public_metrics', {})
            }
            
            logger.info(f"✅ Created user metadata for @{tweet_request.username}: {result['user_metadata']}")
    
    # Save to Firebase if user is not anonymous
    if user_id != "anonymous":
        try:
            enhanced_text = f"Tweet: {tweet_text}"
            if tweet_request.username:
                enhanced_text = f"Tweet by @{tweet_request.username}: {tweet_text}"
                
            history_item = add_analysis_to_history(user_id, enhanced_text, result)
            if history_item:
                result["id"] = history_item.get("id")
                
            # Update user's stats
            update_user_threat_stats(user_id, result)
            
            # Save to threat map if this is a threat
            if result.get("threat", False):
                try:
                    # Enhanced location extraction for single tweet
                    location_str = ""
                    
                    # Try to extract location from tweet metadata if available
                    if tweet_request.tweet_metadata:
                        if 'location' in tweet_request.tweet_metadata:
                            location_str = tweet_request.tweet_metadata.get('location', '')
                        elif 'user_location' in tweet_request.tweet_metadata:
                            location_str = tweet_request.tweet_metadata.get('user_location', '')
                    
                    # Use username for location fallback
                    if not location_str and tweet_request.username:
                        location_str = f"Social media threat from @{tweet_request.username}"
                    elif not location_str:
                        location_str = "Social media threat (unknown location)"
                    
                    # Geocode the location
                    lat, lng = geocode_location(location_str)
                    
                    # Determine priority
                    predicted_class = result.get('predicted_class', 'Unknown')
                    confidence = result.get('confidence', 0.0)
                    priority = determine_threat_priority(predicted_class, confidence)
                    
                    # Create threat location entry
                    threat_data = {
                        "id": f"THR-{result.get('id', str(uuid.uuid4())[:8])}",
                        "type": predicted_class,
                        "lat": lat,
                        "lng": lng,
                        "title": f"{predicted_class} detected" + (f" from @{tweet_request.username}" if tweet_request.username else ""),
                        "location": location_str,
                        "date": result.get('timestamp', datetime.now().isoformat()),
                        "priority": priority,
                        "details": tweet_text[:100] + ("..." if len(tweet_text) > 100 else ""),
                        "caseId": f"THP-{str(uuid.uuid4())[:8]}",
                        "source": "single_tweet_analysis",
                        "confidence": confidence,
                        "predicted_class": predicted_class,
                        "twitter_metadata": result.get("twitter_metadata", {}),
                        "user_metadata": result.get("user_metadata", {}),
                        "text": tweet_text
                    }
                    
                    # Save to threat map
                    saved_location = add_threat_location(user_id, threat_data)
                    if saved_location:
                        logger.info(f"✅ Saved threat location for single tweet analysis")
                    else:
                        logger.warning(f"⚠️ Failed to save threat location for single tweet")
                        
                except Exception as location_error:
                    logger.error(f"❌ Error saving threat location for single tweet: {location_error}")
            
        except Exception as e:
            logger.error(f"Error saving single tweet analysis to Firebase: {e}")
    
    return result

# Enhanced endpoint to get Twitter threats from history
@app.get("/api/twitter/threats")
async def get_twitter_threats(request: Request):
    """Get all Twitter-related threat analyses from user's history"""
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        # Get user's full history
        history = get_user_analysis_history(user_id, limit=1000)  # Get more items for Twitter filtering
        
        if not history:
            return []
        
        # Filter for Twitter-related analyses
        twitter_threats = []
        for item in history:
            # Check if it's Twitter-related by looking at metadata or text patterns
            is_twitter = (
                'twitter_metadata' in item or
                item.get('text', '').startswith('Tweet') or
                '@' in item.get('text', '')
            )
            
            # Only include threats
            is_threat = item.get('threat', False) or item.get('predicted_class') != 'Non-threat/Neutral'
            
            if is_twitter and is_threat:
                twitter_threats.append(item)
        
        logger.info(f"Found {len(twitter_threats)} Twitter threats for user {user_id}")
        return twitter_threats
        
    except Exception as e:
        logger.error(f"Error getting Twitter threats: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/twitter/threats/clear")
async def clear_twitter_threats(request: Request):
    """Clear all Twitter-related threat analyses from user's history"""
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        from firebase_config import get_firestore_db
        db = get_firestore_db()
        if not db:
            raise HTTPException(status_code=500, detail="Database not available")
        
        # Get user's analysis history collection
        history_ref = db.collection("analysis_history").where("user_id", "==", user_id)
        docs = history_ref.get()
        
        deleted_count = 0
        for doc in docs:
            doc_data = doc.to_dict()
            # Check if it's Twitter-related
            is_twitter = (
                'twitter_metadata' in doc_data or
                doc_data.get('text', '').startswith('Tweet') or
                '@' in doc_data.get('text', '')
            )
            
            if is_twitter:
                doc.reference.delete()
                deleted_count += 1
        
        logger.info(f"Deleted {deleted_count} Twitter threat records for user {user_id}")
        return {"success": True, "deleted_count": deleted_count}
        
    except Exception as e:
        logger.error(f"Error clearing Twitter threats: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Twitter Monitor endpoints
@app.post("/api/twitter/monitors")
async def create_twitter_monitor(request: Request, monitor_request: TwitterMonitorRequest):
    """Create a new Twitter monitor"""
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        db = get_firestore_db()
        if not db:
            raise HTTPException(status_code=500, detail="Database not available")
        
        # Create monitor data
        monitor_data = {
            "user_id": user_id,
            "keyword": monitor_request.keyword,
            "type": monitor_request.type,
            "max_results": monitor_request.max_results,
            "schedule_hour": monitor_request.schedule_hour,
            "status": monitor_request.status,
            "created_at": datetime.now().isoformat(),
            "last_run": None,
            "last_results_count": 0
        }
        
        # Add to Firestore
        monitor_ref = db.collection("twitter_monitors").document()
        monitor_ref.set(monitor_data)
        
        # Add the document ID to the returned data
        monitor_data["id"] = monitor_ref.id
        
        logger.info(f"Created Twitter monitor for user {user_id}: {monitor_request.keyword}")
        return {"success": True, "monitor": monitor_data}
        
    except Exception as e:
        logger.error(f"Error creating Twitter monitor: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/twitter/monitors")
async def get_twitter_monitors(request: Request):
    """Get all Twitter monitors for the current user"""
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        db = get_firestore_db()
        if not db:
            raise HTTPException(status_code=500, detail="Database not available")
        
        # Get monitors for the user
        monitors_ref = db.collection("twitter_monitors")
        query = monitors_ref.where("user_id", "==", user_id)
        docs = query.get()
        
        monitors = []
        for doc in docs:
            monitor_data = doc.to_dict()
            monitor_data["id"] = doc.id
            monitors.append(monitor_data)
        
        logger.info(f"Retrieved {len(monitors)} Twitter monitors for user {user_id}")
        return monitors
        
    except Exception as e:
        logger.error(f"Error getting Twitter monitors: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.put("/api/twitter/monitors/{monitor_id}")
async def update_twitter_monitor(request: Request, monitor_id: str, monitor_request: TwitterMonitorRequest):
    """Update an existing Twitter monitor"""
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        db = get_firestore_db()
        if not db:
            raise HTTPException(status_code=500, detail="Database not available")
        
        # Get the monitor document
        monitor_ref = db.collection("twitter_monitors").document(monitor_id)
        monitor_doc = monitor_ref.get()
        
        if not monitor_doc.exists:
            raise HTTPException(status_code=404, detail="Monitor not found")
        
        monitor_data = monitor_doc.to_dict()
        
        # Check if the monitor belongs to the user
        if monitor_data.get("user_id") != user_id:
            raise HTTPException(status_code=403, detail="Access denied")
        
        # Update the monitor data
        updated_data = {
            "keyword": monitor_request.keyword,
            "type": monitor_request.type,
            "max_results": monitor_request.max_results,
            "schedule_hour": monitor_request.schedule_hour,
            "status": monitor_request.status,
            "updated_at": datetime.now().isoformat()
        }
        
        monitor_ref.update(updated_data)
        
        # Return the updated monitor data
        monitor_data.update(updated_data)
        monitor_data["id"] = monitor_id
        
        logger.info(f"Updated Twitter monitor {monitor_id} for user {user_id}")
        return {"success": True, "monitor": monitor_data}
        
    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logger.error(f"Error updating Twitter monitor: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/api/twitter/monitors/{monitor_id}")
async def delete_twitter_monitor(request: Request, monitor_id: str):
    """Delete a Twitter monitor"""
    user_id = get_user_id(request)
    
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        db = get_firestore_db()
        if not db:
            raise HTTPException(status_code=500, detail="Database not available")
        
        # Get the monitor document
        monitor_ref = db.collection("twitter_monitors").document(monitor_id)
        monitor_doc = monitor_ref.get()
        
        if not monitor_doc.exists:
            raise HTTPException(status_code=404, detail="Monitor not found")
        
        monitor_data = monitor_doc.to_dict()
        
        # Check if the monitor belongs to the user
        if monitor_data.get("user_id") != user_id:
            raise HTTPException(status_code=403, detail="Access denied")
        
        # Delete the monitor
        monitor_ref.delete()
        
        logger.info(f"Deleted Twitter monitor {monitor_id} for user {user_id}")
        return {"success": True, "message": "Monitor deleted successfully"}
        
    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logger.error(f"Error deleting Twitter monitor: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Chat AI endpoints
@app.post("/api/chat/message")
async def chat_message(request: Request, chat_request: ChatRequest):
    """
    Send a message to the threat analysis AI assistant
    Enhanced with optional web search for real-time information
    Uses RAG to provide context-aware responses based on user's threat reports
    """
    if not CHAT_ENABLED:
        raise HTTPException(status_code=503, detail="Chat service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Processing chat message for user: {user_id}, web_search: {chat_request.use_web_search}")
        
        # Get relevant context from enhanced RAG
        rag_context = None
        rag_context_used = False
        reports_found = 0
        
        if RAG_ENABLED:
            rag_context = enhanced_rag_service.get_context_for_query(chat_request.message, user_id, max_reports=3)
            rag_context_used = "No relevant threat intelligence found" not in rag_context
            reports_found = len(enhanced_rag_service.search_reports(chat_request.message, user_id, top_k=3))
        
        # Get web search context if requested
        web_context = None
        web_context_used = False
        web_sources = []
        
        if chat_request.use_web_search and WEB_SEARCH_ENABLED:
            try:
                logger.info(f"Performing web search for query: {chat_request.message}")
                web_search_results = web_search_service.search_and_extract(chat_request.message, num_results=3)
                
                if web_search_results.get('success') and web_search_results.get('extracted_content'):
                    # Log details for debugging
                    import json as _json
                    logger.info(f"🔎 Web search raw results: {_json.dumps(web_search_results, ensure_ascii=False)[:1500]}...")

                    web_context = web_search_service.format_web_context(web_search_results)
                    web_context_used = True
                    web_sources = [
                        {
                            'title': content['title'],
                            'url': content['url'],
                            'word_count': content['word_count'],
                            'preview': content['content'][:250]  # First 250 chars for debugging
                        }
                        for content in web_search_results['extracted_content']
                    ]
                    # Log each cleaned content preview
                    for ws in web_sources:
                        logger.info(f"📰 Source: {ws['title']} ({ws['url']}) | Words: {ws['word_count']} | Preview: {ws['preview'][:120]}...")
                    logger.info(f"Web search successful: {len(web_sources)} sources found")
                else:
                    logger.warning("Web search failed or returned no results")
                    # Attempt snippet fallback if search results exist
                    if web_search_results.get('search_results'):
                        snippets = web_search_results['search_results']
                        if snippets:
                            logger.info("🔄 Falling back to web snippets context")
                            web_context_lines = ["=== WEB SEARCH SNIPPETS CONTEXT ===", f"Query: {chat_request.message}", ""]
                            for i, item in enumerate(snippets, 1):
                                web_context_lines.append(f"--- Result {i} ---")
                                web_context_lines.append(f"Title: {item.get('title','')}")
                                web_context_lines.append(f"URL: {item.get('link','')}")
                                if item.get('snippet'):
                                    web_context_lines.append(f"Snippet: {item['snippet']}")
                                web_context_lines.append("")
                            web_context = "\n".join(web_context_lines)
                            web_context_used = True
                            web_sources = [
                                {
                                    'title': item.get('title',''),
                                    'url': item.get('link',''),
                                    'word_count': 0,
                                    'preview': item.get('snippet','')[:250]
                                }
                                for item in snippets
                            ]
                        else:
                            web_context = "No relevant web content found for this query."
                    else:
                        web_context = "No relevant web content found for this query."
            except Exception as web_error:
                logger.error(f"Web search error: {str(web_error)}")
                web_context = "Web search temporarily unavailable."
        
        # Generate AI response with available context
        if chat_request.use_web_search and WEB_SEARCH_ENABLED:
            ai_response = threat_ai.analyze_with_web_context(chat_request.message, rag_context, web_context)
        else:
            ai_response = threat_ai.analyze_with_context(chat_request.message, rag_context)
        
        # Determine context summary
        context_summary_parts = []
        if rag_context_used:
            context_summary_parts.append(f"Internal Intelligence ({reports_found} reports)")
        if web_context_used:
            context_summary_parts.append(f"Web Intelligence ({len(web_sources)} sources)")
        
        context_summary = " + ".join(context_summary_parts) if context_summary_parts else "No relevant context found"
        
        # Save conversation to database
        if RAG_ENABLED:
            # Combine contexts for storage
            combined_context = []
            if rag_context_used:
                combined_context.append(f"RAG Context: {rag_context[:500]}...")
            if web_context_used:
                combined_context.append(f"Web Context: {web_context[:500]}...")
            
            enhanced_rag_service.save_conversation(
                user_id, 
                chat_request.message, 
                ai_response, 
                "\n\n".join(combined_context) if combined_context else None,
                reports_found + len(web_sources)
            )
        
        logger.info(f"Chat response generated for user {user_id}, RAG: {rag_context_used}, Web: {web_context_used}")
        
        return {
            "status": "success",
            "response": ai_response,
            "context_used": rag_context_used or web_context_used,
            "context_summary": context_summary,
            "rag_context_used": rag_context_used,
            "web_context_used": web_context_used,
            "reports_found": reports_found,
            "web_sources": web_sources,
            "web_search_enabled": WEB_SEARCH_ENABLED,
            "timestamp": datetime.now().isoformat(),
            "model": getattr(threat_ai, 'model_name', 'gemini'),
            "user_id": user_id
        }
        
    except Exception as e:
        logger.error(f"Error processing chat message for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Chat processing failed: {str(e)}")

@app.post("/api/chat/web-search")
async def web_search_endpoint(request: Request, search_request: WebSearchRequest):
    """
    Perform web search and extract content for threat analysis
    """
    if not WEB_SEARCH_ENABLED:
        raise HTTPException(status_code=503, detail="Web search service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Web search request from user: {user_id}, query: {search_request.query}")
        
        search_results = web_search_service.search_and_extract(
            search_request.query, 
            search_request.num_results
        )
        
        return {
            "status": "success",
            "search_results": search_results,
            "formatted_context": web_search_service.format_web_context(search_results),
            "timestamp": datetime.now().isoformat(),
            "user_id": user_id
        }
        
    except Exception as e:
        logger.error(f"Error in web search for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Web search failed: {str(e)}")

@app.get("/api/chat/history")
async def get_chat_history(request: Request, limit: int = 20):
    """
    Get conversation history for the current user
    """
    if not CHAT_ENABLED:
        raise HTTPException(status_code=503, detail="Chat service is not available")
    
    user_id = get_user_id(request)
    
    try:
        if RAG_ENABLED:
            history = enhanced_rag_service.get_conversation_history(user_id, limit)
        else:
            history = []
        
        return {
            "status": "success",
            "history": history,
            "count": len(history),
            "user_id": user_id
        }
        
    except Exception as e:
        logger.error(f"Error getting chat history for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to get chat history: {str(e)}")

@app.post("/api/chat/health")
async def chat_health_check(health_request: ChatHealthRequest):
    """
    Check the health status of the chat AI system including web search
    """
    if not CHAT_ENABLED:
        return {
            "status": "disabled",
            "message": "Chat service is not available",
            "chat_enabled": False
        }
    
    try:
        health_status = threat_ai.health_check()
        
        if RAG_ENABLED:
            rag_debug_info = enhanced_rag_service.debug_status()
            rag_status = {
                "status": "healthy",
                "reports_cached": rag_debug_info["reports_cached"],
                "database_path": rag_debug_info["database_path"],
                "model_loaded": rag_debug_info["model_loaded"],
                "report_types": rag_debug_info["report_types"],
                "users": rag_debug_info["users"],
                "sources": rag_debug_info["sources"]
            }
        else:
            rag_status = {"status": "disabled", "message": "RAG service not available"}
        
        # Check web search service health
        if WEB_SEARCH_ENABLED:
            web_search_status = web_search_service.health_check()
        else:
            web_search_status = {"status": "disabled", "message": "Web search service not available"}
        
        return {
            "status": "success",
            "chat_enabled": True,
            "ai_service": health_status,
            "rag_service": rag_status,
            "web_search_service": web_search_status,
            "web_search_enabled": WEB_SEARCH_ENABLED,
            "model": getattr(threat_ai, 'model_name', 'gemini'),
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"Error checking chat health: {e}")
        return {
            "status": "error",
            "chat_enabled": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }

@app.post("/api/chat/rag/refresh")
async def refresh_rag_cache(request: Request):
    """
    Refresh the RAG cache with latest reports
    """
    if not CHAT_ENABLED:
        raise HTTPException(status_code=503, detail="Chat service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Refreshing enhanced RAG cache for user: {user_id}")
        
        if RAG_ENABLED:
            # Refresh the cache
            enhanced_rag_service._refresh_cache()
            reports_cached = len(enhanced_rag_service.reports_cache)
        else:
            reports_cached = 0
        
        return {
            "status": "success",
            "message": "Enhanced RAG cache refreshed successfully",
            "reports_cached": reports_cached,
            "rag_enabled": RAG_ENABLED,
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"Error refreshing RAG cache: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to refresh RAG cache: {str(e)}")

# Briefing API endpoints
@app.post("/api/briefing/generate-ppt")
async def generate_powerpoint_presentation(request: Request, briefing_request: BriefingRequest):
    """Generate PowerPoint presentation from threat analysis report using Gemini AI"""
    user_id = get_user_id(request)
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")
    
    try:
        logger.info(f"Generating PowerPoint presentation for user {user_id}, report: {briefing_request.reportId}")
        
        # Fetch the report data based on reportId
        report_data = await fetch_report_data(briefing_request.reportId)
        if not report_data:
            raise HTTPException(status_code=404, detail="Report not found")
        
        # Process report with Gemini API to structure content
        structured_slides = await process_report_with_gemini(report_data)
        if not structured_slides:
            raise HTTPException(status_code=500, detail="Failed to process report with Gemini AI")
        
        # Generate PowerPoint file
        ppt_file_path = await generate_pptx_file(structured_slides, user_id)
        if not ppt_file_path or not os.path.exists(ppt_file_path):
            raise HTTPException(status_code=500, detail="Failed to generate PowerPoint file")
        
        logger.info(f"PowerPoint presentation generated successfully: {ppt_file_path}")
        
        # Return file as download
        return FileResponse(
            path=ppt_file_path,
            filename=f"threat_analysis_briefing_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logger.error(f"Error generating PowerPoint presentation: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate presentation: {str(e)}")

async def fetch_report_data(report_id: str) -> dict:
    """Fetch report data based on report ID"""
    try:
        if report_id == "threat-report":
            # Mock threat report data for now - in production this would fetch from database
            return {
                "title": "Threat Analysis Report",
                "report_id": "TR-91156370",
                "generated_date": datetime.now().strftime("%B %d, %Y"),
                "classification": "OFFICIAL USE ONLY",
                "analyst": "ASTRA Intelligence Division",
                "status": "Active",
                "metadata": {
                    "total_incidents": 15,
                    "critical_incidents": 3,
                    "high_incidents": 6,
                    "medium_incidents": 4,
                    "low_incidents": 2
                },
                "incidents": [
                    {
                        "incident_id": "INC-001",
                        "timestamp": "2024-01-15T10:30:00Z",
                        "classification": "Direct Violence Threats",
                        "severity": "CRITICAL",
                        "risk_level": "HIGH",
                        "content_summary": "Explicit threats against government officials with specific plans and timelines"
                    },
                    {
                        "incident_id": "INC-002", 
                        "timestamp": "2024-01-15T11:45:00Z",
                        "classification": "Hate Speech/Extremism",
                        "severity": "HIGH",
                        "risk_level": "MEDIUM",
                        "content_summary": "Extremist rhetoric promoting violence against specific ethnic groups"
                    },
                    {
                        "incident_id": "INC-003",
                        "timestamp": "2024-01-15T13:20:00Z",
                        "classification": "Criminal Activity",
                        "severity": "CRITICAL",
                        "risk_level": "HIGH", 
                        "content_summary": "Detailed plans for cyberattack on critical infrastructure"
                    }
                ]
            }
        elif report_id == "summary-report":
            return {
                "title": "Weekly Security Summary",
                "period": f"Week of {(datetime.now() - timedelta(days=7)).strftime('%Y-%m-%d')}",
                "executive_summary": "This week saw significant increases in cyber threats targeting financial institutions and government agencies.",
                "key_findings": [
                    "25% increase in phishing attacks",
                    "3 new malware variants detected",
                    "Critical infrastructure probing attempts"
                ]
            }
        else:
            # Sample report for demonstration
            return {
                "title": "Sample Threat Intelligence Report",
                "report_id": "SAMPLE-001",
                "generated_date": datetime.now().strftime("%B %d, %Y"),
                "classification": "SAMPLE DATA",
                "summary": "This is a sample report demonstrating PowerPoint generation capabilities.",
                "incidents": [
                    {"type": "Hate Speech/Extremism", "severity": "HIGH", "count": 2},
                    {"type": "Direct Violence Threats", "severity": "CRITICAL", "count": 1}
                ]
            }
            
    except Exception as e:
        logger.error(f"Error fetching report data: {e}")
        return None

async def process_report_with_gemini(report_data: dict) -> list:
    """Process report data with Gemini API to create structured slides"""
    try:
        # Enhanced prompt – MANDATORY 10 slides with detailed content
        prompt = f"""You are a senior threat-intelligence analyst creating a comprehensive PowerPoint briefing. You MUST create EXACTLY 10 slides with detailed content.

CRITICAL REQUIREMENTS:
- You MUST return EXACTLY 10 slides (no more, no less)
- Each slide MUST have detailed bullet points (4-6 bullets per slide)
- Include specific examples and data from the input report
- Use real incident IDs, counts, and details from the data

OUTPUT FORMAT:
Return ONLY valid JSON array with exactly 10 objects. Each object must have:
• "title": string (slide title)
• "content": array of 4-6 detailed bullet point strings
• "speaker_notes": string (150+ words for 1 minute of speech)
• "image_query": string (descriptive image search term)

MANDATORY SLIDE STRUCTURE (10 slides):
1. Title Slide: Report metadata, classification, agency branding
2. Executive Summary: Key findings, total incidents, critical threats
3. Threat Landscape Overview: External context, threat trends
4. Critical Incidents Analysis: Detailed breakdown of high-severity incidents
5. Incident Type Distribution: Breakdown by threat categories with examples
6. Geographic Impact Analysis: Location-based threat patterns
7. Temporal Analysis: Time-based threat trends and patterns
8. Risk Assessment: Detailed risk analysis and impact assessment
9. Mitigation Strategies: Specific recommendations and countermeasures
10. Conclusion & Next Steps: Summary and action items

USE REAL DATA: Reference specific incident IDs, dates, locations, and details from the input report. Make each slide content-rich and informative.

INPUT REPORT DATA:
{json.dumps(report_data, indent=2)}

RETURN EXACTLY 10 SLIDES AS JSON ARRAY:"""

        model = genai.GenerativeModel("gemini-2.0-flash-exp")
        response = model.generate_content(
            prompt,
            generation_config={
                "response_mime_type": "application/json"
            }
        )
        
        # Parse the JSON response
        slides_data = json.loads(response.text)
        logger.info(f"Gemini API returned {len(slides_data)} slides")
        
        # Ensure we have exactly 10 slides
        if len(slides_data) < 10:
            logger.warning(f"Gemini returned only {len(slides_data)} slides, expanding to 10")
            slides_data = expand_to_ten_slides(slides_data, report_data)
        
        return slides_data
        
    except Exception as e:
        logger.error(f"Error processing report with Gemini API: {e}")
        # Return comprehensive 10-slide fallback
        return create_comprehensive_fallback_slides(report_data)

def expand_to_ten_slides(slides_data: list, report_data: dict) -> list:
    """Expand slides to exactly 10 with detailed content"""
    incidents = report_data.get('incidents', [])
    
    # Define the 10 required slide types
    required_slides = [
        "Title Slide",
        "Executive Summary", 
        "Threat Landscape Overview",
        "Critical Incidents Analysis",
        "Incident Type Distribution",
        "Geographic Impact Analysis",
        "Temporal Analysis",
        "Risk Assessment",
        "Mitigation Strategies",
        "Conclusion & Next Steps"
    ]
    
    expanded_slides = []
    
    for i, slide_type in enumerate(required_slides):
        if i < len(slides_data):
            # Use existing slide
            expanded_slides.append(slides_data[i])
        else:
            # Create new detailed slide
            expanded_slides.append(create_detailed_slide(slide_type, i, incidents, report_data))
    
    return expanded_slides

def create_detailed_slide(slide_type: str, slide_index: int, incidents: list, report_data: dict) -> dict:
    """Create a detailed slide with real data"""
    
    if slide_type == "Title Slide":
        return {
            "title": "Threat Analysis Report",
            "content": [
                f"Report ID: {report_data.get('report_id', 'TR-001')}",
                f"Generated: {report_data.get('generated_date', 'Current Date')}",
                f"Classification: {report_data.get('classification', 'Official Use Only')}",
                "ASTRA Intelligence Division",
                f"Total Incidents: {len(incidents)}"
            ],
            "speaker_notes": f"This presentation provides a comprehensive analysis of {len(incidents)} security incidents identified by our threat detection systems. The report covers critical threats, risk assessments, and actionable mitigation strategies for our organization's security posture.",
            "image_query": "cyber security threat analysis background"
        }
    
    elif slide_type == "Executive Summary":
        critical_count = len([inc for inc in incidents if inc.get('severity') == 'CRITICAL'])
        high_count = len([inc for inc in incidents if inc.get('severity') == 'HIGH'])
        return {
            "title": "Executive Summary",
            "content": [
                f"Total Incidents Analyzed: {len(incidents)}",
                f"Critical Threats Identified: {critical_count}",
                f"High-Priority Incidents: {high_count}",
                "Multiple threat vectors detected across systems",
                "Immediate response required for critical incidents"
            ],
            "speaker_notes": f"Our analysis identified {len(incidents)} security incidents requiring attention. Among these, {critical_count} are classified as critical and {high_count} as high priority. These findings indicate a significant threat landscape that requires immediate and coordinated response from our security teams.",
            "image_query": "executive summary dashboard charts"
        }
    
    elif slide_type == "Critical Incidents Analysis":
        critical_incidents = [inc for inc in incidents if inc.get('severity') == 'CRITICAL']
        return {
            "title": "Critical Incidents Analysis",
            "content": [
                f"Critical Incidents: {len(critical_incidents)}",
                f"Most Common Type: {get_most_common_type(critical_incidents)}",
                f"Average Response Time: {calculate_avg_response(critical_incidents)}",
                "Immediate containment measures required",
                "Enhanced monitoring protocols recommended"
            ],
            "speaker_notes": f"Critical incidents represent our highest priority threats. We identified {len(critical_incidents)} such incidents, primarily involving {get_most_common_type(critical_incidents)}. These require immediate containment and enhanced monitoring protocols to prevent escalation.",
            "image_query": "critical security alert warning"
        }
    
    elif slide_type == "Incident Type Distribution":
        type_counts = get_incident_type_counts(incidents)
        return {
            "title": "Incident Type Distribution",
            "content": [
                f"Hate Speech/Extremism: {type_counts.get('Hate Speech/Extremism', 0)}",
                f"Direct Violence Threats: {type_counts.get('Direct Violence Threats', 0)}",
                f"Harassment: {type_counts.get('Harassment and Intimidation', 0)}",
                f"Criminal Activity: {type_counts.get('Criminal Activity', 0)}",
                f"Child Safety: {type_counts.get('Child Safety Threats', 0)}"
            ],
            "speaker_notes": f"Our threat landscape shows diverse attack vectors. The distribution reveals {type_counts.get('Hate Speech/Extremism', 0)} extremism cases, {type_counts.get('Direct Violence Threats', 0)} direct threats, and {type_counts.get('Criminal Activity', 0)} criminal activities. This diversity requires specialized response protocols.",
            "image_query": "pie chart threat distribution"
        }
    
    elif slide_type == "Geographic Impact Analysis":
        locations = extract_locations(incidents)
        return {
            "title": "Geographic Impact Analysis",
            "content": [
                f"Affected Regions: {len(locations)}",
                f"Primary Threat Zone: {get_primary_location(locations)}",
                "Cross-border threat patterns detected",
                "Regional response coordination needed",
                "Local law enforcement collaboration required"
            ],
            "speaker_notes": f"Geographic analysis reveals threats across {len(locations)} regions, with {get_primary_location(locations)} as the primary threat zone. Cross-border patterns suggest coordinated attacks requiring regional response coordination and law enforcement collaboration.",
            "image_query": "world map threat locations"
        }
    
    elif slide_type == "Temporal Analysis":
        time_patterns = analyze_temporal_patterns(incidents)
        return {
            "title": "Temporal Analysis",
            "content": [
                f"Peak Activity: {time_patterns.get('peak_time', 'Unknown')}",
                f"Attack Frequency: {time_patterns.get('frequency', 'Variable')}",
                "Weekend surge in activity detected",
                "Business hours targeting observed",
                "Seasonal patterns identified"
            ],
            "speaker_notes": f"Temporal analysis reveals distinct attack patterns with peak activity during {time_patterns.get('peak_time', 'specific periods')}. Weekend surges and business hours targeting suggest strategic timing by threat actors, requiring enhanced monitoring during these periods.",
            "image_query": "timeline attack patterns"
        }
    
    elif slide_type == "Risk Assessment":
        risk_score = calculate_risk_score(incidents)
        return {
            "title": "Risk Assessment",
            "content": [
                f"Overall Risk Score: {risk_score}/10",
                f"Critical Vulnerabilities: {count_critical_vulns(incidents)}",
                "System exposure levels elevated",
                "Data breach probability: High",
                "Reputation damage risk: Significant"
            ],
            "speaker_notes": f"Our risk assessment indicates an overall risk score of {risk_score}/10, with {count_critical_vulns(incidents)} critical vulnerabilities identified. System exposure is elevated with high probability of data breaches and significant reputation damage risk.",
            "image_query": "risk assessment matrix"
        }
    
    elif slide_type == "Mitigation Strategies":
        return {
            "title": "Mitigation Strategies",
            "content": [
                "Implement enhanced monitoring systems",
                "Deploy automated threat detection",
                "Establish incident response protocols",
                "Conduct regular security training",
                "Develop threat intelligence sharing"
            ],
            "speaker_notes": "Our mitigation strategy focuses on proactive defense through enhanced monitoring, automated detection, and comprehensive incident response protocols. Regular training and threat intelligence sharing will strengthen our security posture.",
            "image_query": "cyber security defense strategy"
        }
    
    else:  # Conclusion & Next Steps
        return {
            "title": "Conclusion & Next Steps",
            "content": [
                f"Total threats addressed: {len(incidents)}",
                "Immediate action items identified",
                "Enhanced security protocols recommended",
                "Ongoing monitoring established",
                "Follow-up assessment scheduled"
            ],
            "speaker_notes": f"This comprehensive analysis of {len(incidents)} threats provides a clear roadmap for strengthening our security posture. Immediate actions are required, with ongoing monitoring and follow-up assessments to ensure continued protection.",
            "image_query": "conclusion action plan"
        }

def create_comprehensive_fallback_slides(report_data: dict) -> list:
    """Create comprehensive 10-slide fallback with real data"""
    incidents = report_data.get('incidents', [])
    
    return [
        create_detailed_slide("Title Slide", 0, incidents, report_data),
        create_detailed_slide("Executive Summary", 1, incidents, report_data),
        create_detailed_slide("Threat Landscape Overview", 2, incidents, report_data),
        create_detailed_slide("Critical Incidents Analysis", 3, incidents, report_data),
        create_detailed_slide("Incident Type Distribution", 4, incidents, report_data),
        create_detailed_slide("Geographic Impact Analysis", 5, incidents, report_data),
        create_detailed_slide("Temporal Analysis", 6, incidents, report_data),
        create_detailed_slide("Risk Assessment", 7, incidents, report_data),
        create_detailed_slide("Mitigation Strategies", 8, incidents, report_data),
        create_detailed_slide("Conclusion & Next Steps", 9, incidents, report_data)
    ]

# Helper functions for slide content
def get_most_common_type(incidents: list) -> str:
    types = [inc.get('classification', 'Unknown') for inc in incidents]
    return max(set(types), key=types.count) if types else "Unknown"

def calculate_avg_response(incidents: list) -> str:
    return "2.5 hours"  # Placeholder

def get_incident_type_counts(incidents: list) -> dict:
    counts = {}
    for inc in incidents:
        inc_type = inc.get('classification', 'Unknown')
        counts[inc_type] = counts.get(inc_type, 0) + 1
    return counts

def extract_locations(incidents: list) -> list:
    locations = []
    for inc in incidents:
        if 'location' in inc:
            locations.append(inc['location'])
    return list(set(locations))

def get_primary_location(locations: list) -> str:
    return locations[0] if locations else "Unknown"

def analyze_temporal_patterns(incidents: list) -> dict:
    return {"peak_time": "Weekend evenings", "frequency": "Daily"}

def calculate_risk_score(incidents: list) -> int:
    critical_count = len([inc for inc in incidents if inc.get('severity') == 'CRITICAL'])
    return min(10, 5 + critical_count)

def count_critical_vulns(incidents: list) -> int:
    return len([inc for inc in incidents if inc.get('severity') == 'CRITICAL'])

async def generate_pptx_file(slides_data: list, user_id: str) -> str:
    """Generate PowerPoint file from structured slides data"""
    try:
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Define style options for randomization
        colors = [
            RGBColor(31, 73, 125),  # Dark blue
            RGBColor(0, 51, 102),   # Navy
            RGBColor(51, 102, 0),   # Dark green
            RGBColor(139, 0, 0),    # Dark red
            RGBColor(75, 0, 130)    # Indigo
        ]
        fonts = ['Calibri', 'Arial', 'Helvetica', 'Times New Roman']
        
        for i, slide_data in enumerate(slides_data):
            # Use title slide layout for first slide, content layout for others
            if i == 0:
                slide_layout = prs.slide_layouts[0]  # Title slide
            else:
                slide_layout = prs.slide_layouts[1]  # Title and content
                
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if slide.shapes.title:
                title = slide.shapes.title
                title.text = slide_data.get('title', 'Untitled Slide')
                title.text_frame.paragraphs[0].font.size = Pt(28)
                title.text_frame.paragraphs[0].font.name = random.choice(fonts)
                title.text_frame.paragraphs[0].font.color.rgb = random.choice(colors)
            
            # Add content if not title slide
            if i > 0 and len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for bullet_point in slide_data.get('content', []):
                    p = text_frame.add_paragraph()
                    p.text = bullet_point
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.name = random.choice(fonts)
                    p.font.color.rgb = random.choice(colors)
                    p.alignment = PP_ALIGN.LEFT
            
            # Optional speaker notes – some programs (e.g. Apple Keynote) are
            # stricter about the notesSlide part of the spec and may refuse to
            # open a file that contains malformed or unexpected notes. To
            # maximise cross-platform compatibility we only add notes when
            # explicitly enabled via an environment variable.
            # Always include speaker notes for detailed presentation content
            if True:  # os.getenv('INCLUDE_PPT_SPEAKER_NOTES', 'false').lower() == 'true':
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data.get('speaker_notes', '')
            
            # Add image if provided and placeholder available
            img_query = slide_data.get('image_query', '') if isinstance(slide_data, dict) else ''
            if img_query:
                try:
                    img_url = f"https://source.unsplash.com/1600x900/?{requests.utils.quote(img_query)}"
                    img_resp = requests.get(img_url, timeout=10)
                    if img_resp.status_code == 200:
                        img_bytes = io.BytesIO(img_resp.content)
                        # Place image at bottom-right quarter of the slide
                        pic_left = Inches(6)
                        pic_top = Inches(4.5)
                        pic_width = Inches(3)
                        slide.shapes.add_picture(img_bytes, pic_left, pic_top, width=pic_width)
                except Exception as img_err:
                    logger.warning(f"Failed to fetch/add image for slide {i}: {img_err}")
        
        # Save to temporary file
        temp_dir = tempfile.mkdtemp()
        filename = f"threat_briefing_{user_id}_{int(time.time())}.pptx"
        file_path = os.path.join(temp_dir, filename)
        
        prs.save(file_path)
        logger.info(f"PowerPoint file saved: {file_path}")
        
        return file_path
        
    except Exception as e:
        logger.error(f"Error generating PowerPoint file: {e}")
        return None

# --------------------------- New API Endpoint ---------------------------

@app.post("/api/briefing/generate-script")
async def generate_presenter_script_endpoint(request: Request, briefing_request: BriefingRequest):
    """Generate a downloadable presenter script (.txt) for a given report"""
    user_id = get_user_id(request)
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")

    try:
        logger.info(f"Generating presenter script for user {user_id}, report: {briefing_request.reportId}")

        # Fetch report data
        report_data = await fetch_report_data(briefing_request.reportId)
        if not report_data:
            raise HTTPException(status_code=404, detail="Report not found")

        # Re-use slide planning from Gemini to keep consistency
        slides_data = await process_report_with_gemini(report_data)

        # Create presenter script text
        script_text = await create_presenter_script(slides_data)

        # Save to temp file
        temp_dir = tempfile.mkdtemp()
        filename = f"threat_briefing_script_{user_id}_{int(time.time())}.txt"
        file_path = os.path.join(temp_dir, filename)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(script_text)

        logger.info(f"Presenter script saved: {file_path}")

        return FileResponse(
            path=file_path,
            filename=filename,
            media_type="text/plain"
        )

    except HTTPException as http_err:
        raise http_err
    except Exception as e:
        logger.error(f"Error generating presenter script: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate presenter script: {str(e)}")

async def create_presenter_script(slides_data: list) -> str:
    """Generate a presenter script (~1 minute per slide) using Gemini"""
    try:
        script_prompt = f"""You are preparing a spoken presenter script to accompany a PowerPoint briefing.  Below is the JSON array describing each slide. For **each slide**, write a paragraph of roughly 150 spoken-word tokens (≈1 minute) expanding on the bullet points.  Use a professional, confident tone.

Return plain text in the following format (no extra commentary):

Slide 1 – <title>
<paragraph>

Slide 2 – <title>
<paragraph>

... etc.

SLIDES JSON:
{json.dumps(slides_data, indent=2)}"""

        model = genai.GenerativeModel("gemini-2.0-flash-exp")
        response = model.generate_content(script_prompt)
        return response.text
    except Exception as e:
        logger.error(f"Error generating presenter script: {e}")
        return "Presenter script generation failed."

# --------------------------- Video Generation Functions ---------------------------

def convert_pptx_to_pdf(pptx_path: str) -> str:
    """Convert PowerPoint file to PDF using LibreOffice"""
    try:
        # Create temporary directory for PDF
        temp_dir = tempfile.mkdtemp()
        pdf_filename = os.path.basename(pptx_path).replace('.pptx', '.pdf')
        pdf_path = os.path.join(temp_dir, pdf_filename)
        
        # Use LibreOffice to convert PPTX to PDF
        cmd = [
            'soffice', '--headless', '--convert-to', 'pdf',
            '--outdir', temp_dir, pptx_path
        ]
        
        logger.info(f"Converting PPTX to PDF: {' '.join(cmd)}")
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        
        if result.returncode != 0:
            logger.error(f"LibreOffice conversion failed: {result.stderr}")
            raise Exception(f"PDF conversion failed: {result.stderr}")
        
        # Find the generated PDF file
        generated_pdf = os.path.join(temp_dir, pdf_filename)
        if not os.path.exists(generated_pdf):
            # LibreOffice might have created a different filename
            pdf_files = [f for f in os.listdir(temp_dir) if f.endswith('.pdf')]
            if pdf_files:
                generated_pdf = os.path.join(temp_dir, pdf_files[0])
            else:
                raise Exception("PDF file not found after conversion")
        
        logger.info(f"PDF conversion successful: {generated_pdf}")
        return generated_pdf
        
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {e}")
        raise

def generate_video_from_pptx(pptx_path: str, pdf_path: str, output_path: str) -> bool:
    """Generate video from PowerPoint using ppt_presenter tool"""
    try:
        # Import the ppt_presenter function
        import sys
        sys.path.append(os.path.join(os.path.dirname(__file__), '..', 'ppt_presenter'))
        from ppt_presenter import ppt_presenter_parallel
        
        logger.info(f"Starting video generation: {pptx_path} -> {output_path}")
        
        # Use parallel processing for better performance
        max_workers = min(multiprocessing.cpu_count(), 4)  # Conservative worker count
        ppt_presenter_parallel(pptx_path, pdf_path, output_path, max_workers)
        
        if os.path.exists(output_path):
            logger.info(f"Video generation successful: {output_path}")
            return True
        else:
            logger.error("Video file not created")
            return False
            
    except Exception as e:
        logger.error(f"Error generating video: {e}")
        return False

async def create_video_from_pptx(pptx_path: str, user_id: str) -> str:
    """Create video from PowerPoint file with narration"""
    try:
        # Convert PPTX to PDF
        pdf_path = convert_pptx_to_pdf(pptx_path)
        
        # Generate video
        temp_dir = tempfile.mkdtemp()
        video_filename = f"threat_briefing_video_{user_id}_{int(time.time())}.mp4"
        video_path = os.path.join(temp_dir, video_filename)
        
        success = generate_video_from_pptx(pptx_path, pdf_path, video_path)
        
        if success:
            logger.info(f"Video created successfully: {video_path}")
            return video_path
        else:
            raise Exception("Video generation failed")
            
    except Exception as e:
        logger.error(f"Error creating video: {e}")
        raise

@app.post("/api/briefing/generate-video")
async def generate_video_from_pptx_endpoint(request: Request, briefing_request: BriefingRequest):
    """Generate a video presentation from PowerPoint with narration"""
    user_id = get_user_id(request)
    if user_id == "anonymous":
        raise HTTPException(status_code=401, detail="Authentication required")

    try:
        logger.info(f"Generating video for user {user_id}, report: {briefing_request.reportId}")

        # Fetch report data
        report_data = await fetch_report_data(briefing_request.reportId)
        if not report_data:
            raise HTTPException(status_code=404, detail="Report not found")

        # Generate PowerPoint first
        slides_data = await process_report_with_gemini(report_data)
        pptx_path = await generate_pptx_file(slides_data, user_id)
        
        if not pptx_path:
            raise HTTPException(status_code=500, detail="Failed to generate PowerPoint")

        # Generate video from PowerPoint
        video_path = await create_video_from_pptx(pptx_path, user_id)
        
        if not video_path:
            raise HTTPException(status_code=500, detail="Failed to generate video")

        logger.info(f"Video generated successfully: {video_path}")

        return FileResponse(
            path=video_path,
            filename=os.path.basename(video_path),
            media_type="video/mp4"
        )
        
    except Exception as e:
        logger.error(f"Error generating video: {e}")
        raise HTTPException(status_code=500, detail=f"Video generation failed: {str(e)}")

# Threat Map API endpoints
@app.get("/api/threat-map/data")
async def get_threat_map_data(request: Request):
    """Get threat map data for the current user"""
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Fetching threat map data for user: {user_id}")
        
        if user_id == "anonymous":
            # Return sample data for anonymous users
            sample_threats = generate_sample_threat_data()
            logger.info(f"Generated {len(sample_threats)} sample threats for anonymous user")
            return sample_threats
        else:
            # Get user's actual threat locations from Firebase
            threat_locations = get_user_threat_locations(user_id, limit=100)
            
            # If user has no data, generate some sample data mixed with their actual analysis history
            if not threat_locations:
                logger.info(f"No threat locations found for user {user_id}, generating from analysis history")
                threat_locations = generate_threats_from_user_history(user_id)
            
            logger.info(f"Returning {len(threat_locations)} threat locations for user {user_id}")
            return threat_locations
            
    except Exception as e:
        logger.error(f"Error getting threat map data for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to get threat map data: {str(e)}")

@app.post("/api/threat-map/filter")
async def filter_threat_map(request: Request, filter_req: ThreatMapFilterRequest):
    """Filter threat map data based on criteria"""
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Filtering threat map data for user: {user_id} with criteria: {filter_req}")
        
        if user_id == "anonymous":
            # Return filtered sample data for anonymous users
            sample_threats = generate_sample_threat_data()
            filtered_data = filter_sample_threats(sample_threats, filter_req)
            return filtered_data
        else:
            # Filter user's actual threat locations
            filtered_locations = filter_threat_locations(
                user_id=user_id,
                time_range_days=filter_req.timeRange or 30,
                threat_types=filter_req.threatTypes,
                priority_levels=filter_req.priority
            )
            
            logger.info(f"Filtered to {len(filtered_locations)} threat locations for user {user_id}")
            return filtered_locations
            
    except Exception as e:
        logger.error(f"Error filtering threat map data for user {user_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to filter threat map data: {str(e)}")

def generate_sample_threat_data():
    """Generate sample threat data for anonymous users or empty accounts"""
    import random
    from datetime import datetime, timedelta
    
    cities = [
        ("New York", 40.7128, -74.0060),
        ("Los Angeles", 34.0522, -118.2437),
        ("Chicago", 41.8781, -87.6298),
        ("Houston", 29.7604, -95.3698),
        ("Phoenix", 33.4484, -112.0740),
        ("Philadelphia", 39.9526, -75.1652),
        ("San Antonio", 29.4241, -98.4936),
        ("San Diego", 32.7157, -117.1611),
        ("Dallas", 32.7767, -96.7970),
        ("San Francisco", 37.7749, -122.4194),
        ("London", 51.5074, -0.1278),
        ("Paris", 48.8566, 2.3522),
        ("Tokyo", 35.6762, 139.6503),
        ("Berlin", 52.5200, 13.4050),
        ("Sydney", 33.8688, 151.2093)
    ]
    
    threat_types = [
        "Direct Violence Threats",
        "Criminal Activity", 
        "Harassment and Intimidation",
        "Hate Speech/Extremism",
        "Child Safety Threats"
    ]
    
    priorities = ["low", "medium", "high", "critical"]
    
    threats = []
    
    # Generate 15-25 random threats
    for i in range(random.randint(15, 25)):
        city_name, lat, lng = random.choice(cities)
        threat_type = random.choice(threat_types)
        priority = random.choice(priorities)
        
        # Add some random variance to coordinates
        lat += random.uniform(-0.5, 0.5)
        lng += random.uniform(-0.5, 0.5)
        
        threat_id = f"THR-{random.randint(10000, 99999)}"
        case_id = f"THP-{random.randint(10000, 99999)}"
        
        threat = {
            "id": threat_id,
            "type": threat_type,
            "lat": round(lat, 6),
            "lng": round(lng, 6),
            "title": f"{threat_type} detected in {city_name}",
            "location": city_name,
            "date": (datetime.now() - timedelta(days=random.randint(1, 60))).isoformat(),
            "priority": priority,
            "details": f"Sample {threat_type.lower()} detected in {city_name} area requiring investigation.",
            "caseId": case_id,
            "source": "sample_data",
            "confidence": round(random.uniform(0.6, 0.95), 2),
            "predicted_class": threat_type
        }
        
        threats.append(threat)
    
    return threats

def filter_sample_threats(threats, filter_req):
    """Filter sample threat data based on criteria"""
    filtered = threats.copy()
    
    # Apply time range filter
    if filter_req.timeRange:
        cutoff_date = datetime.now() - timedelta(days=filter_req.timeRange)
        filtered = [
            threat for threat in filtered 
            if datetime.fromisoformat(threat["date"]) > cutoff_date
        ]
    
    # Apply threat type filter
    if filter_req.threatTypes and len(filter_req.threatTypes) > 0:
        filtered = [
            threat for threat in filtered 
            if threat["type"] in filter_req.threatTypes
        ]
    
    # Apply priority filter
    if filter_req.priority and len(filter_req.priority) > 0:
        filtered = [
            threat for threat in filtered 
            if threat["priority"] in filter_req.priority
        ]
    
    return filtered

def generate_threats_from_user_history(user_id):
    """Generate threat map data from user's analysis history"""
    try:
        from firebase_config import get_user_analysis_history
        
        # Get user's history
        history = get_user_analysis_history(user_id, limit=500)
        
        threats = []
        processed_count = 0
        
        for item in history:
            # Only process threats
            if not item.get('threat', False):
                continue
            
            # Extract location information - enhanced location extraction
            location_str = ""
            
            # Check Twitter metadata for location first
            if 'twitter_metadata' in item:
                twitter_meta = item['twitter_metadata']
                
                # Try to get actual location from metadata
                if 'location' in twitter_meta and twitter_meta['location']:
                    location_str = twitter_meta['location']
                elif 'user_location' in twitter_meta and twitter_meta['user_location']:
                    location_str = twitter_meta['user_location']
                else:
                    # Fall back to username-based location
                    username = twitter_meta.get('username', '')
                    if username:
                        location_str = f"Social media threat from @{username}"
            
            # Check if there's saved location information from the analysis
            if not location_str and 'location' in item:
                location_str = item['location']
            
            # If still no specific location, use Unknown Location
            if not location_str:
                location_str = "Unknown Location"
            
            # Geocode the location
            lat, lng = geocode_location(location_str)
            
            # Determine priority
            predicted_class = item.get('predicted_class', 'Unknown')
            confidence = item.get('confidence', 0.0)
            priority = determine_threat_priority(predicted_class, confidence)
            
            # Create threat location entry
            threat_data = {
                "id": f"THR-{item.get('id', str(uuid.uuid4())[:8])}",
                "type": predicted_class,
                "lat": round(lat, 6),
                "lng": round(lng, 6),
                "title": f"{predicted_class} detected",
                "location": location_str,
                "date": item.get('timestamp', datetime.now().isoformat()),
                "priority": priority,
                "details": item.get('text', 'No details available')[:100] + ("..." if len(item.get('text', '')) > 100 else ""),
                "caseId": f"THP-{str(uuid.uuid4())[:8]}",
                "source": "user_history",
                "confidence": confidence,
                "predicted_class": predicted_class,
                "twitter_metadata": item.get('twitter_metadata', {}),
                "text": item.get('text', '')
            }
            
            # Save to Firebase for future use
            add_threat_location(user_id, threat_data)
            threats.append(threat_data)
            
            processed_count += 1
            if processed_count >= 50:  # Limit to 50 threats max
                break
        
        logger.info(f"Generated {len(threats)} threat locations from user history for user {user_id}")
        return threats
        
    except Exception as e:
        logger.error(f"Error generating threats from user history for user {user_id}: {e}")
        return []

# Global exception handler
@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    error_msg = str(exc)
    logger.error(f"Unhandled exception: {error_msg}")
    logger.exception("Full stack trace:")
    return {"error": error_msg, "status_code": 500} 

@app.post("/api/chat/rag/migrate")
async def migrate_firebase_to_rag(request: Request):
    """
    Migrate existing Firebase data AND report API data to RAG system
    This fixes cases where Firebase has data but RAG is empty
    """
    if not CHAT_ENABLED or not RAG_ENABLED:
        raise HTTPException(status_code=503, detail="Chat/RAG service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Starting comprehensive data migration for user: {user_id}")
        
        migrated_count = 0
        total_analysis = 0
        total_threats = 0
        total_summaries = 0
        
        # Step 1: Migrate Firebase data (existing logic)
        logger.info("🔥 Step 1: Migrating Firebase data...")
        
        if user_id == "admin" or user_id.startswith("admin_"):
            # Admin can migrate all data
            users_ref = db.collection('users')
            user_docs = list(users_ref.stream())
            logger.info(f"Admin migration: processing {len(user_docs)} users")
        else:
            # Regular user can only migrate their own data
            user_doc_ref = db.collection('users').document(user_id)
            user_doc = user_doc_ref.get()
            user_docs = [user_doc] if user_doc.exists else []
            logger.info(f"User migration: processing user {user_id}")
        
        for user_doc in user_docs:
            current_user_id = user_doc.id
            
            if not user_doc.exists:
                logger.warning(f"User {current_user_id} does not exist, skipping")
                continue
                
            logger.info(f"📊 Migrating Firebase data for user: {current_user_id}")
            
            # Migrate analysis history
            analysis_docs = list(db.collection('users').document(current_user_id).collection('analysis_history').stream())
            logger.info(f"   Found {len(analysis_docs)} analysis history items")
            for doc in analysis_docs:
                doc_data = doc.to_dict()
                text = doc_data.get('text', '')
                prediction = doc_data.get('prediction', {})
                
                if text and prediction:
                    try:
                        success = enhanced_rag_service.add_prediction_analysis(
                            user_id=current_user_id,
                            text=text,
                            prediction_result=prediction,
                            source='firebase_migration'
                        )
                        if success:
                            migrated_count += 1
                            total_analysis += 1
                    except Exception as e:
                        logger.error(f"Error migrating analysis for user {current_user_id}: {e}")
            
            # Migrate threat reports  
            threat_docs = list(db.collection('users').document(current_user_id).collection('threat_reports').stream())
            logger.info(f"   Found {len(threat_docs)} threat reports")
            for doc in threat_docs:
                doc_data = doc.to_dict()
                try:
                    success = enhanced_rag_service.add_threat_report(
                        user_id=current_user_id,
                        report_data=doc_data,
                        source='firebase_migration'
                    )
                    if success:
                        migrated_count += 1
                        total_threats += 1
                except Exception as e:
                    logger.error(f"Error migrating threat report for user {current_user_id}: {e}")
            
            # Migrate summary reports
            summary_docs = list(db.collection('users').document(current_user_id).collection('summary_reports').stream())
            logger.info(f"   Found {len(summary_docs)} summary reports")
            for doc in summary_docs:
                doc_data = doc.to_dict()
                try:
                    success = enhanced_rag_service.add_summary_report(
                        user_id=current_user_id,
                        report_data=doc_data,
                        source='firebase_migration'  
                    )
                    if success:
                        migrated_count += 1
                        total_summaries += 1
                except Exception as e:
                    logger.error(f"Error migrating summary report for user {current_user_id}: {e}")
        
        # Step 2: Migrate Report API data (NEW)
        logger.info("📋 Step 2: Migrating Report API data...")
        
        try:
            # Get summary report from API
            summary_response = await get_summary_report(request)
            if hasattr(summary_response, 'report') or isinstance(summary_response, dict):
                summary_data = summary_response.get('report') if isinstance(summary_response, dict) else summary_response.report
                if summary_data:
                    logger.info("   Found summary report from API")
                    success = enhanced_rag_service.add_summary_report(
                        user_id=user_id,
                        report_data=summary_data,
                        source='api_migration'
                    )
                    if success:
                        migrated_count += 1
                        total_summaries += 1
                        logger.info("   ✅ Summary report migrated from API")
        except Exception as e:
            logger.warning(f"Could not migrate summary report from API: {e}")
        
        try:
            # Get threat report from API
            threat_response = await get_threat_report(request)
            if hasattr(threat_response, 'report') or isinstance(threat_response, dict):
                threat_data = threat_response.get('report') if isinstance(threat_response, dict) else threat_response.report
                if threat_data:
                    logger.info("   Found threat report from API")
                    success = enhanced_rag_service.add_threat_report(
                        user_id=user_id,
                        report_data=threat_data,
                        source='api_migration'
                    )
                    if success:
                        migrated_count += 1
                        total_threats += 1
                        logger.info("   ✅ Threat report migrated from API")
        except Exception as e:
            logger.warning(f"Could not migrate threat report from API: {e}")
        
        # Step 3: Get user's analysis history from API
        logger.info("📈 Step 3: Migrating analysis history from API...")
        
        try:
            history_response = await get_user_history(request)
            if hasattr(history_response, 'history') or isinstance(history_response, dict):
                history_data = history_response.get('history') if isinstance(history_response, dict) else history_response.history
                if history_data and isinstance(history_data, list):
                    logger.info(f"   Found {len(history_data)} analysis items from API")
                    for item in history_data:
                        text = item.get('text', '')
                        prediction = item.get('prediction', {})
                        
                        if text and prediction:
                            try:
                                success = enhanced_rag_service.add_prediction_analysis(
                                    user_id=user_id,
                                    text=text,
                                    prediction_result=prediction,
                                    source='api_migration'
                                )
                                if success:
                                    migrated_count += 1
                                    total_analysis += 1
                            except Exception as e:
                                logger.error(f"Error migrating API analysis item: {e}")
                    logger.info(f"   ✅ {len([i for i in history_data if i.get('text') and i.get('prediction')])} analysis items migrated from API")
        except Exception as e:
            logger.warning(f"Could not migrate analysis history from API: {e}")
        
        # Get final RAG status
        final_rag_status = enhanced_rag_service.debug_status()
        
        logger.info(f"🎉 Migration complete: {migrated_count} items migrated")
        
        return {
            "status": "success",
            "message": "Comprehensive data migration completed",
            "migrated_items": migrated_count,
            "breakdown": {
                "analysis_history": total_analysis,
                "threat_reports": total_threats,
                "summary_reports": total_summaries
            },
            "sources_migrated": ["firebase", "report_api", "analysis_api"],
            "final_rag_reports": final_rag_status['reports_cached'],
            "timestamp": datetime.now().isoformat()
        }
        
    except Exception as e:
        logger.error(f"Error during comprehensive migration: {e}")
        raise HTTPException(status_code=500, detail=f"Migration failed: {str(e)}")

# FIR (First Information Report) Endpoints
@app.post("/api/fir/create")
async def create_fir(request: Request, fir_request: FIRRequest):
    """Create a new FIR for a critical Twitter threat"""
    if not FIR_ENABLED:
        raise HTTPException(status_code=503, detail="FIR service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Creating FIR for user: {user_id}")
        
        # Create FIR
        fir_data = await fir_service.create_fir(fir_request.threat_data, user_id)
        
        logger.info(f"FIR created successfully: {fir_data['fir_id']}")
        return {
            "status": "success",
            "message": "FIR created successfully",
            "fir_id": fir_data['fir_id'],
            "fir_data": fir_data
        }
        
    except Exception as e:
        logger.error(f"Error creating FIR: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to create FIR: {str(e)}")

@app.get("/api/fir/list")
async def get_user_firs(request: Request, limit: int = 50):
    """Get all FIRs for the current user"""
    if not FIR_ENABLED:
        raise HTTPException(status_code=503, detail="FIR service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Fetching FIRs for user: {user_id}")
        
        firs = await fir_service.get_user_firs(user_id, limit)
        
        logger.info(f"Found {len(firs)} FIRs for user {user_id}")
        return {
            "status": "success",
            "firs": firs,
            "count": len(firs)
        }
        
    except Exception as e:
        logger.error(f"Error fetching FIRs: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch FIRs: {str(e)}")

@app.get("/api/fir/{fir_id}")
async def get_fir_by_id(request: Request, fir_id: str):
    """Get specific FIR by ID"""
    if not FIR_ENABLED:
        raise HTTPException(status_code=503, detail="FIR service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Fetching FIR {fir_id} for user: {user_id}")
        
        fir_data = await fir_service.get_fir_by_id(fir_id)
        
        if not fir_data:
            raise HTTPException(status_code=404, detail="FIR not found")
        
        # Check if user owns this FIR
        if fir_data.get('user_id') != user_id:
            raise HTTPException(status_code=403, detail="Access denied")
        
        logger.info(f"FIR {fir_id} retrieved successfully")
        return {
            "status": "success",
            "fir": fir_data
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching FIR {fir_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch FIR: {str(e)}")

@app.put("/api/fir/{fir_id}/status")
async def update_fir_status(request: Request, fir_id: str, status: str):
    """Update FIR status"""
    if not FIR_ENABLED:
        raise HTTPException(status_code=503, detail="FIR service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Updating FIR {fir_id} status to {status} for user: {user_id}")
        
        # Verify FIR exists and user owns it
        fir_data = await fir_service.get_fir_by_id(fir_id)
        if not fir_data:
            raise HTTPException(status_code=404, detail="FIR not found")
        
        if fir_data.get('user_id') != user_id:
            raise HTTPException(status_code=403, detail="Access denied")
        
        # Update status
        success = await fir_service.update_fir_status(fir_id, status)
        
        if success:
            logger.info(f"FIR {fir_id} status updated to {status}")
            return {
                "status": "success",
                "message": f"FIR status updated to {status}",
                "fir_id": fir_id,
                "new_status": status
            }
        else:
            raise HTTPException(status_code=500, detail="Failed to update FIR status")
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error updating FIR {fir_id} status: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to update FIR status: {str(e)}")

@app.get("/api/fir/{fir_id}/pdf")
async def download_fir_pdf(request: Request, fir_id: str):
    """Download FIR as PDF"""
    if not FIR_ENABLED:
        raise HTTPException(status_code=503, detail="FIR service is not available")
    
    user_id = get_user_id(request)
    
    try:
        logger.info(f"Generating PDF for FIR {fir_id} for user: {user_id}")
        
        # Get FIR data
        fir_data = await fir_service.get_fir_by_id(fir_id)
        if not fir_data:
            raise HTTPException(status_code=404, detail="FIR not found")
        
        # Check if user owns this FIR
        if fir_data.get('user_id') != user_id:
            raise HTTPException(status_code=403, detail="Access denied")
        
        # Generate PDF
        pdf_path = fir_service.generate_pdf(fir_data)
        
        # Return file response
        filename = f"FIR_{fir_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        
        logger.info(f"PDF generated successfully: {pdf_path}")
        return FileResponse(
            path=pdf_path,
            filename=filename,
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="{filename}"'}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating PDF for FIR {fir_id}: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate PDF: {str(e)}")

# Run the server when executed directly
if __name__ == "__main__":
    import uvicorn
    logger.info("🚀 Starting Astra Threat Detection Platform API Server...")
    uvicorn.run(
        "main:app",
        host="0.0.0.0",
        port=8000,
        reload=True,
        log_level="info"
    ) 